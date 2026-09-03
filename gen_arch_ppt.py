# -*- coding: utf-8 -*-
"""
生成《方案 B2 架构图集》PPTX —— 全部使用 PPT 原生矩形 / 线条 / 文本绘制，不嵌入任何图片。
同时输出 /tmp/arch_ops.json 供 PIL 预览渲染检查版式。
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

# ---------------------------------------------------------------- palette ---
INK    = '0F172A'; SUB = '475569'; MUT = '64748B'; FAINT = '94A3B8'
BORDER = 'E2E8F0'; GRID = 'EEF2F7'; BG = 'F8FAFC'
AZ     = '0078D4'; AZD = '1E40AF'          # azure blue
GRN    = '059669'; GRND = '047857'         # scheme B green
PUR    = '7C3AED'; PURD = '6D28D9'         # databricks / read
CYN    = '0891B2'; CYND = '0E7490'         # app layer cyan
ORG    = 'EA580C'                          # ingest orange
AMB    = 'D97706'; AMBD = '92400E'         # orchestration amber
ROS    = 'E11D48'                          # retention rose
SLA    = '475569'                          # slate

L_BLUE  = 'EFF6FF'; B_BLUE  = 'BFDBFE'
L_GRN   = 'ECFDF5'; B_GRN   = 'A7F3D0'
L_PUR   = 'F5F3FF'; B_PUR   = 'DDD6FE'
L_CYN   = 'ECFEFF'; B_CYN   = 'A5F3FC'
L_AMB   = 'FFFBEB'; B_AMB   = 'FDE68A'
L_ROS   = 'FFF1F2'; B_ROS   = 'FECDD3'
L_VNET  = 'F0F7FC'
L_GRNB  = 'F0FDF4'; B_GRNB  = '86EFAC'
L_PURB  = 'F5F3FF'; B_PURB  = 'C4B5FD'
L_CYNB  = 'ECFEFF'; B_CYNB  = '67E8F9'

FONT_LAT = 'Segoe UI'; FONT_EA = 'Microsoft YaHei'; FONT_MONO = 'Consolas'
OPS = []          # preview ops
_cur = {'i': 0}
def _slide_idx(): return _cur['i']

# ------------------------------------------------------------------ helpers -
def _rgb(h): return RGBColor.from_string(h)

def _mk_shadow(shape, blur=55000, dist=22000, alpha=12000, color=INK):
    spPr = shape._element.spPr
    el = parse_xml(
        '<a:effectLst %s><a:outerShdw blurRad="%d" dist="%d" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:outerShdw></a:effectLst>'
        % (nsdecls('a'), blur, dist, color, alpha))
    spPr.append(el)

def _no_line_shadow(shape):
    try: shape.shadow.inherit = False
    except Exception: pass

def rect(sl, x, y, w, h, fill=None, line_c=None, lw=1.0, rad=None, dash=None,
         shadow=False, shape_type=None):
    st = shape_type or (MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE)
    sp = sl.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if rad is not None:
        try: sp.adjustments[0] = max(0.0, min(0.5, rad / min(w, h)))
        except Exception: pass
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    else: sp.fill.background()
    if line_c:
        sp.line.color.rgb = _rgb(line_c); sp.line.width = Pt(lw)
        if dash: sp.line.dash_style = dash
    else:
        sp.line.fill.background()
    if shadow: _mk_shadow(sp)
    else: _no_line_shadow(sp)
    OPS.append({'t': 'r', 's': _slide_idx(), 'x': x, 'y': y, 'w': w, 'h': h,
                'fill': fill, 'line': line_c, 'lw': lw,
                'dash': dash.name if hasattr(dash, 'name') else dash, 'rad': rad})
    return sp

def oval(sl, x, y, w, h, fill=None, line_c=None, lw=1.0, dash=None):
    return rect(sl, x, y, w, h, fill, line_c, lw, None, dash, False, MSO_SHAPE.OVAL)

def _ln_style(conn, color, w, dash, arrow, head=False):
    conn.line.color.rgb = _rgb(color); conn.line.width = Pt(w)
    if dash: conn.line.dash_style = dash
    ln = conn.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    # 关键：显式清空效果列表，阻断主题 effectRef 的默认软阴影（否则箭头会"发虚/模糊"）
    spPr = conn._element.spPr
    if spPr.find(qn('a:effectLst')) is None:
        spPr.append(spPr.makeelement(qn('a:effectLst'), {}))
    # 移除主题样式引用，彻底避免任何主题默认效果/线型干扰
    style = conn._element.find(qn('p:style'))
    if style is not None:
        conn._element.remove(style)

def seg(sl, x1, y1, x2, y2, color=MUT, w=1.0, dash=None, arrow=True, head=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _ln_style(c, color, w, dash, arrow, head)
    OPS.append({'t': 'l', 's': _slide_idx(), 'x1': x1, 'y1': y1, 'x2': x2, 'y2': y2,
                'c': color, 'w': w, 'dash': dash.name if hasattr(dash, 'name') else dash, 'ar': arrow})

def poly(sl, pts, color=MUT, w=1.0, dash=None, arrow=True):
    for i in range(len(pts) - 1):
        last = (i == len(pts) - 2)
        seg(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], color, w, dash, arrow and last)

ALIGN = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}
ANCHOR = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}

def _apply_run(r, text, sz, b, c, mono=False, spc=None, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(sz); f.bold = b; f.italic = italic
    f.color.rgb = _rgb(c)
    f.name = FONT_MONO if mono else FONT_LAT
    rPr = r._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_EA})
    if latin is not None: latin.addnext(ea)
    else: rPr.append(ea)
    if spc: rPr.set('spc', str(spc))

def txt(sl, x, y, w, h, paras, sz=8, c=INK, b=False, align='l', anchor='t',
        leading=1.12, mono=False, spc=None, vert=False, rec=True):
    """paras: str | list; item = str | (str, style-dict) | dict(runs=[(s, st)], align=, sb=, ls=)"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try: tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception: pass
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = ANCHOR[anchor]
    if vert: tf._txBody.bodyPr.set('vert', 'eaVert')
    if isinstance(paras, str): paras = [paras]
    first = True
    for item in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ALIGN[align]
        p.line_spacing = leading
        p.space_before = Pt(0); p.space_after = Pt(0)
        sb = 0; it_runs = []
        if isinstance(item, str):
            it_runs = [(item, {})]
        elif isinstance(item, tuple):
            it_runs = [item]
        elif isinstance(item, dict):
            it_runs = item.get('runs', [])
            if 'align' in item: p.alignment = ALIGN[item['align']]
            if 'ls' in item: p.line_spacing = item['ls']
            sb = item.get('sb', 0)
        if sb: p.space_before = Pt(sb)
        for s, st in it_runs:
            r = p.add_run()
            _apply_run(r, s, st.get('sz', sz), st.get('b', b), st.get('c', c),
                       st.get('mono', mono), st.get('spc', spc), st.get('i', False))
    if rec:
        OPS.append({'t': 'x', 's': _slide_idx(), 'x': x, 'y': y, 'w': w, 'h': h,
                    'paras': [_norm_para(it, sz, b, c) for it in paras],
                    'align': align, 'anchor': anchor, 'leading': leading, 'vert': vert})
    return tb

def _norm_para(item, sz, b, c):
    if isinstance(item, str): return {'runs': [(item, {'sz': sz, 'b': b, 'c': c})]}
    if isinstance(item, tuple): return {'runs': [(item[0], _merge(item[1], sz, b, c))]}
    if isinstance(item, dict):
        return {'runs': [(s, _merge(st, sz, b, c)) for s, st in item.get('runs', [])],
                'align': item.get('align'), 'sb': item.get('sb', 0)}
    return {'runs': []}

def _merge(st, sz, b, c):
    d = {'sz': sz, 'b': b, 'c': c}; d.update(st or {}); return d

def chip(sl, x, y, w, h, title, sub=None, fill='FFFFFF', line_c=BORDER, lw=1.0, rad=0.05,
         tc=INK, tsz=8, sc=MUT, ssz=6.5, dash=None, shadow=False, align='c', sub2=None,
         anchor='m', leading=1.08):
    rect(sl, x, y, w, h, fill, line_c, lw, rad, dash, shadow)
    if isinstance(title, str):
        paras = [(title, {})]
    elif isinstance(title, dict):
        paras = [title]
    else:
        paras = list(title)
    if sub: paras.append((sub, {'sz': ssz, 'c': sc}))
    if sub2: paras.append((sub2, {'sz': ssz, 'c': sc}))
    txt(sl, x + 0.03, y, w - 0.06, h, paras, sz=tsz, c=tc, align=align, anchor=anchor,
        leading=leading)

# ------------------------------------------------------------------- icons -
def person(sl, cx, cy, color=AZ, s=1.0):
    d = 0.066 * s
    oval(sl, cx - d / 2, cy - 0.085 * s, d, d, color)
    bw = 0.115 * s; bh = 0.075 * s
    rect(sl, cx - bw / 2, cy - 0.008 * s, bw, bh, color, rad=0.03 * s)

def dbicon(sl, cx, cy, color=SLA, s=1.0):
    w = 0.13 * s; hh = 0.028 * s; body = 0.062 * s
    rect(sl, cx - w / 2, cy - body / 2 - hh / 4, w, body, color)
    oval(sl, cx - w / 2, cy - body / 2 - hh, w, hh, color)
    oval(sl, cx - w / 2, cy + body / 2 - hh / 2 + hh / 4, w, hh, 'FFFFFF', line_c=color, lw=0.75)

def cloud(sl, cx, cy, color=AZ, s=1.0):
    oval(sl, cx - 0.052 * s, cy - 0.030 * s, 0.104 * s, 0.075 * s, color)
    oval(sl, cx - 0.075 * s, cy - 0.010 * s, 0.070 * s, 0.052 * s, color)
    oval(sl, cx + 0.005 * s, cy - 0.018 * s, 0.066 * s, 0.058 * s, color)

def lockicon(sl, cx, cy, color=AZ, s=1.0):
    bw = 0.10 * s; bh = 0.075 * s
    rect(sl, cx - bw / 2, cy, bw, bh, color, rad=0.015 * s)
    OPS.append({'t': 'arc', 's': _slide_idx(), 'cx': cx, 'cy': cy, 'r': bw * 0.32, 'c': color, 'w': 1.2})

# ------------------------------------------------------------ slide chrome -
def title_bar(sl, title, right_note, accent, idx_note=''):
    rect(sl, 0.22, 0.17, 0.055, 0.40, accent, rad=0.027)
    txt(sl, 0.38, 0.13, 9.0, 0.34, title, sz=18.5, b=True, c=INK, anchor='m')
    if right_note:
        txt(sl, 9.40, 0.06, 3.05, 0.42, right_note, sz=7.5, c=MUT, align='r', anchor='m', leading=1.25)
    if idx_note:
        txt(sl, 12.52, 0.40, 0.6, 0.18, idx_note, sz=7, c=FAINT, align='r', mono=True)

# ================================================================= COVER ====
def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 0
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb('FFFFFF')
    # decorative bands
    rect(sl, 0, 0, 13.333, 0.12, AZ)
    rect(sl, 0, 0.12, 13.333, 0.035, GRN)
    rect(sl, 0, 0.155, 13.333, 0.035, PUR)
    rect(sl, 7.75, 1.30, 5.58, 5.9, 'F8FAFC', rad=0.25)
    rect(sl, 8.30, 1.85, 4.75, 4.75, 'FFFFFF', line_c=BORDER, lw=1.2, rad=0.22, shadow=True)
    # mini diagram motif (abstract)
    rect(sl, 8.68, 2.42, 4.0, 0.52, L_GRN, B_GRN, 1.2, 0.08)
    txt(sl, 8.68, 2.42, 4.0, 0.52, '访问入口 · 查询 / 审计 / 运维', sz=8.5, c=GRND, b=True, align='c', anchor='m')
    seg(sl, 10.68, 2.94, 10.68, 3.29, GRN, 1.6)
    rect(sl, 8.68, 3.29, 4.0, 0.52, L_CYN, B_CYN, 1.2, 0.08)
    txt(sl, 8.68, 3.29, 4.0, 0.52, '自研应用层 · AKS（React + .NET）', sz=8.5, c=CYND, b=True, align='c', anchor='m')
    seg(sl, 9.55, 3.81, 9.55, 4.13, GRN, 1.4); seg(sl, 11.81, 3.81, 11.81, 4.13, PUR, 1.4)
    rect(sl, 8.68, 4.13, 1.95, 0.52, L_PUR, B_PUR, 1.2, 0.08)
    txt(sl, 8.68, 4.13, 1.95, 0.52, '写侧 · Databricks', sz=8, c=PURD, b=True, align='c', anchor='m')
    rect(sl, 10.73, 4.13, 1.95, 0.52, L_PUR, B_PUR, 1.2, 0.08)
    txt(sl, 10.73, 4.13, 1.95, 0.52, '读侧 · Serverless', sz=8, c=PURD, b=True, align='c', anchor='m')
    seg(sl, 10.68, 4.65, 10.68, 4.99, MUT, 1.4)
    rect(sl, 8.68, 4.99, 4.0, 0.52, L_BLUE, B_BLUE, 1.2, 0.08)
    txt(sl, 8.68, 4.99, 4.0, 0.52, 'ADLS Gen2 · Iceberg 湖仓 + Blob 附件', sz=8, c=AZD, b=True, align='c', anchor='m')
    txt(sl, 8.68, 5.72, 4.0, 0.3, '存算分离 · 私有端点 · 合规销毁', sz=7.5, c=FAINT, align='c')

    txt(sl, 0.9, 1.55, 6.6, 0.3, 'RIMS 退役归档平台 · 方案 B2', sz=11, c=MUT, spc=300, b=True)
    txt(sl, 0.88, 1.92, 6.9, 0.85, '湖仓一体架构图集', sz=33, b=True, c=INK)
    rect(sl, 0.92, 2.86, 0.62, 0.045, AZ)
    txt(sl, 0.92, 3.02, 6.4, 0.3, '架构图 · 数据加密设计 · 元数据库 ER 图 · 全部以 PPT 原生形状绘制',
        sz=10.5, c=MUT)
    items = [
        ('01', AZ, L_BLUE, B_BLUE, 'Azure 基础设施部署架构 · 方案 B2', 'VNet · AKS · Databricks Serverless · NCC 私有端点 · 六大 PaaS 服务'),
        ('02', GRN, L_GRN, B_GRN, '应用架构图 · 方案 B2（湖仓一体 · Databricks SQL Serverless 读）', '访问入口 → 自研应用层（AKS）→ 写 / 读 / 附件三条链路 → 数据底座'),
        ('03', PUR, L_PUR, B_PUR, '方案 B2 · 湖仓一体 · 存算分离架构', '源系统 → 接入配置 → 抽取 → Iceberg 冷湖 → Databricks 读写一体计算层 → 消费入口 + 销毁回路'),
        ('04', AZ, L_BLUE, B_BLUE, '方案 B2 · 数据加密设计', '四层加密模型 · aes_encrypt + UC 列掩码 · 信封加密 / Crypto-Shredding'),
        ('05', SLA, 'F6F8FA', 'CBD5E1', '元数据库 ER 图（35 表 / 11 模块 · 上下两页）', '上篇：平台 / 权限 / 数据源配置　下篇：同步归档 / 动态查询 / 保留销毁'),
    ]
    y = 3.24
    for no, col, lf, bl, t, d in items:
        rect(sl, 0.92, y, 6.35, 0.66, 'FFFFFF', BORDER, 1.2, 0.09, shadow=True)
        rect(sl, 1.06, y + 0.09, 0.48, 0.48, lf, bl, 1.2, 0.08)
        txt(sl, 1.06, y + 0.09, 0.48, 0.48, no, sz=12, b=True, c=col, align='c', anchor='m', mono=True)
        txt(sl, 1.68, y + 0.08, 5.5, 0.24, t, sz=10, b=True, c=INK)
        txt(sl, 1.68, y + 0.35, 5.5, 0.22, d, sz=7.2, c=MUT)
        y += 0.74
    rect(sl, 0.92, 6.98, 6.35, 0.035, BORDER)
    txt(sl, 0.92, 7.08, 6.4, 0.25, '整理自 report/Architecture.html · docs/数据库ER图与全量表结构.md', sz=8, c=FAINT)

# ======================================================== SLIDE 1 · AZURE ===
def slide_azure(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 1
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb(BG)
    title_bar(sl, 'Azure 基础设施部署架构 · 方案 B2', '', AZ)

    # ---- 技术栈 ----
    rect(sl, 0.22, 0.66, 2.02, 1.94, 'FFFFFF', BORDER, 1.2, 0.07)
    txt(sl, 0.36, 0.76, 1.8, 0.2, '技术栈', sz=9, b=True, c=INK)
    rect(sl, 0.36, 0.99, 0.30, 0.025, AZ)
    stack = [
        ('1', '前端 · 配置化查询平台 React', False),
        ('2', '后端 · .NET 10（AKS 容器化）', False),
        ('3', '抽取 · Apache SeaTunnel', False),
        ('4', '编排 · Azure Data Factory', False),
        ('5', '表格式 · Apache Iceberg', False),
        ('6', '引擎 · Databricks Serverless（读写）', False),
        ('7', '治理 · Unity Catalog', False),
        ('8', '附件迁移 · Azure Storage Mover', True),
    ]
    yy = 1.10
    for no, s, hot in stack:
        txt(sl, 0.36, yy, 1.78, 0.17,
            [dict(runs=[(no + '  ', {'b': True, 'c': AMB if hot else AZ, 'sz': 6.8}),
                        (s, {'c': AMBD if hot else SUB, 'sz': 6.8, 'b': hot})])])
        yy += 0.183

    # ---- 终端用户 ----
    rect(sl, 0.22, 2.72, 2.02, 1.98, L_GRN, B_GRN, 1.2, 0.07)
    txt(sl, 0.36, 2.80, 1.8, 0.2, '终端用户', sz=9, b=True, c=GRND)
    users = [('业务查询员', '查本系统归档数据'), ('审计 / 合规岗', '跨系统检索与举证'), ('平台运维', '监控 · 销毁作业值守')]
    yy = 3.06
    for t, s in users:
        chip(sl, 0.34, yy, 1.78, 0.42, t, s, 'FFFFFF', B_GRN, 1.0, 0.05, tc=INK, tsz=7.6, sc=MUT, ssz=6.2)
        person(sl, 0.52, yy + 0.21, AZ, 0.9)
        yy += 0.50
    txt(sl, 0.30, 4.50, 1.86, 0.16, '浏览器 HTTPS 443 · Entra ID 单点登录', sz=6.2, b=True, c=GRN, align='c')

    # ---- 源系统 ----
    rect(sl, 2.38, 0.66, 6.22, 1.60, 'FFFFFF', 'CBD5E1', 1.2, 0.07)
    txt(sl, 2.54, 0.74, 4.4, 0.2, '源系统 · 15 套待退役业务系统', sz=9, b=True, c=SUB)
    txt(sl, 5.6, 0.76, 2.86, 0.16, '退役窗口内只读 · 数据只出不进', sz=6.2, c=FAINT, align='r')
    dbs = [('Oracle', 'TCP 1521'), ('MySQL', 'TCP 3306'), ('DB2 / 其他', 'TCP 50000')]
    xx = 2.54
    for t, s in dbs:
        chip(sl, xx, 1.00, 1.86, 0.50, t, s, 'FFFFFF', 'CBD5E1', 1.0, 0.05, tsz=7.8, sc=MUT, ssz=6.3)
        dbicon(sl, xx + 0.24, 1.25, SLA, 0.95)
        xx += 2.00
    chip(sl, 2.54, 1.62, 2.90, 0.50,
         [dict(runs=[('结构化抽取 · ', {'c': INK, 'sz': 7.4, 'b': True}),
                     ('Apache SeaTunnel', {'c': ORG, 'sz': 7.4, 'b': True})])],
         '仅读权限 · 按勾选的表抽取（10 选 5）· ADF 编排',
         'FFFFFF', ORG, 1.3, 0.05, sc=MUT, ssz=6.0)
    chip(sl, 5.56, 1.62, 2.84, 0.50,
         [dict(runs=[('附件迁移 · ', {'c': AMBD, 'sz': 7.4, 'b': True}),
                     ('Storage Mover', {'c': AMB, 'sz': 7.4, 'b': True})])],
         '读 NFS / S3 / Blob 附件源 · 服务免费',
         L_AMB, AMB, 1.4, 0.05, sc='78350F', ssz=6.0)
    txt(sl, 2.54, 2.10, 5.9, 0.14, '跨机房链路：ExpressRoute / 站点到站点 VPN · 源库防火墙放行 NAT 固定出口 IP',
        sz=6.2, b=True, c=AMB, align='c')

    # ---- 网络连通 mini ----
    rect(sl, 8.74, 0.66, 1.74, 1.60, 'F8FAFC', 'CBD5E1', 1.0, 0.07, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 8.86, 0.74, 1.5, 0.18, '网络连通', sz=8, b=True, c=SUB)
    for i, s in enumerate(['ExpressRoute / S2S VPN 专线',
                           'NAT 固定出口 IP 白名单',
                           'PaaS 全私有端点 · 禁公网',
                           'WAF + App Gateway 收敛入口']):
        txt(sl, 8.86, 0.98 + i * 0.29, 1.52, 0.28,
            [dict(runs=[('· ', {'c': AZ, 'b': True, 'sz': 6.4}), (s, {'c': SUB, 'sz': 6.4})], ls=1.05)])

    # ---- Databricks 托管平面 ----
    rect(sl, 10.60, 0.66, 2.52, 4.04, 'FAF9FF', PUR, 1.4, 0.08, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 10.74, 0.76, 2.3, 0.2, 'Databricks 托管平面', sz=9.5, b=True, c=PUR)
    txt(sl, 10.74, 0.97, 2.3, 0.15, '★ 不在你的 VNet 内 · 运行于 Databricks 账户', sz=6.2, b=True, c='DC2626')
    db_cards = [
        (1.18, 0.74, 'FFFFFF', 'A78BFA', 1.0, '控制平面 Control Plane',
         [('Web UI · 作业调度 · 无状态', False), ('Unity Catalog 元数据服务', False)], PURD, 7.6),
        (1.99, 0.74, L_PUR, '8B5CF6', 1.5, 'Databricks SQL Serverless',
         [('读侧查询 · X-Small 6 DBU/h', False), ('空闲自动缩到 0 · 零运维', True)], PURD, 7.6),
        (2.80, 0.74, 'FFFFFF', 'A78BFA', 1.0, 'Databricks Serverless',
         [('写侧 ETL：RAW→CURATED→LAKE→SERVE', False), ('Compaction · VACUUM · 到期销毁', False)], PURD, 7.6),
        (3.61, 0.60, 'ECFDF5', '6EE7B7', 1.3, 'NCC 网络连接 + 托管私有端点',
         [('Serverless 私有访问 ADLS 的唯一通道', True)], GRND, 7.2),
    ]
    for y0, h0, f, bl, lw0, t, lines, tc, tsz0 in db_cards:
        rect(sl, 10.74, y0, 2.24, h0, f, bl, lw0, 0.05)
        cloud(sl, 10.92, y0 + 0.14, PUR, 0.75)
        txt(sl, 10.80, y0 + 0.06, 2.12, 0.15, t, sz=tsz0, b=True, c=tc, align='c')
        yy = y0 + 0.28
        for s, strong in lines:
            txt(sl, 10.80, yy, 2.12, 0.15, s, sz=6.2, c=GRN if (strong and f == 'ECFDF5') else MUT,
                b=strong, align='c')
            yy += 0.175
    txt(sl, 10.74, 4.34, 2.24, 0.26, 'Premium 档（新建唯一档）· NCC 需 Account Admin', sz=6.0, c=PUR,
        b=True, align='c', anchor='m')

    # ---- Azure VNet ----
    rect(sl, 2.38, 2.42, 8.10, 2.28, L_VNET, AZ, 1.5, 0.08)
    txt(sl, 2.54, 2.50, 3.4, 0.2, 'Azure VNet · 10.0.0.0/16', sz=9.5, b=True, c=AZ)
    txt(sl, 6.3, 2.54, 3.6, 0.15, 'Hub-Spoke · NSG 管控 · PaaS 全部关闭公网', sz=6.4, c=MUT, align='r')
    # edge row
    chip(sl, 2.54, 2.74, 0.96, 0.46, 'WAF', 'OWASP 防护', 'FFFFFF', AZ, 1.2, 0.05, tsz=7.8, sc=MUT, ssz=6.0)
    chip(sl, 3.66, 2.74, 1.46, 0.46, 'Application Gateway', 'TLS 终止 · 七层路由', 'FFFFFF', AZ, 1.2, 0.05, tsz=7.4, sc=MUT, ssz=6.0)
    chip(sl, 5.28, 2.74, 1.06, 0.46, 'NAT Gateway', '出站固定 IP', 'FFFFFF', AZ, 1.2, 0.05, tsz=7.4, sc=MUT, ssz=6.0)
    seg(sl, 3.50, 2.97, 3.64, 2.97, AZ, 1.4)
    seg(sl, 5.12, 2.97, 5.26, 2.97, AZ, 1.4)
    seg(sl, 4.39, 3.20, 4.39, 3.32, AZ, 1.4)   # App Gateway -> AKS
    txt(sl, 6.75, 2.80, 3.55, 0.34,
        [dict(runs=[('私有 DNS：', {'b': True, 'c': SUB, 'sz': 6.4}),
                    ('privatelink.* 区域集中解析，AKS / Databricks 经内网域名访问 PaaS', {'c': MUT, 'sz': 6.4})], ls=1.15)])

    # AKS subnet
    rect(sl, 2.54, 3.34, 4.60, 1.24, 'FFFFFF', AZ, 1.1, 0.06)
    txt(sl, 2.66, 3.42, 2.6, 0.17, 'AKS 子网 · 10.0.1.0/24', sz=7.8, b=True, c=AZ)
    txt(sl, 4.6, 3.44, 2.4, 0.14, '跨 AZ 节点池 · HPA · 托管标识免密', sz=5.9, c=FAINT, align='r')
    rect(sl, 2.66, 3.62, 4.36, 0.56, None, 'CBD5E1', 1.0, 0.05, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 2.74, 3.65, 2.2, 0.12, 'Node · 跨可用区冗余 · 托管节点池', sz=5.6, b=True, c=SUB)
    chip(sl, 2.76, 3.78, 2.10, 0.36, '前端 Pod · React', '配置化查询界面 · JSON Schema',
         L_BLUE, AZ, 1.4, 0.05, tc=AZD, tsz=7.2, sc=MUT, ssz=5.9)
    chip(sl, 4.94, 3.78, 2.00, 0.36, '后端 Pod · .NET 单体', '元数据 · 权限 · 查询代理 · 审计 · 销毁',
         'F0F9FF', CYND, 1.4, 0.05, tc=CYND, tsz=7.2, sc=MUT, ssz=5.9)
    rect(sl, 2.66, 4.26, 4.36, 0.24, 'F0F9FF', '7DD3FC', 1.0, 0.04)
    txt(sl, 2.72, 4.26, 4.24, 0.24, '元数据服务：系统台账 · 表清单勾选 · 保留期与销毁策略 · 任务记录 · 操作审计',
        sz=5.9, b=True, c='0C4A6E', align='c', anchor='m')

    # PE subnet
    rect(sl, 7.30, 3.34, 3.06, 1.24, 'FFFFFF', AZ, 1.1, 0.06)
    txt(sl, 7.42, 3.42, 2.9, 0.17, '私有端点子网 · 10.0.2.0/24', sz=7.8, b=True, c=AZ)
    pes = ['ADLS 私有端点（blob + dfs）', 'SQL Database 私有端点', 'Key Vault 私有端点', 'Databricks 工作区端点']
    for i, s in enumerate(pes):
        px = 7.42 + (i % 2) * 1.46; py = 3.64 + (i // 2) * 0.40
        chip(sl, px, py, 1.40, 0.34, s, None, L_BLUE, B_BLUE, 1.0, 0.04, tc=AZD, tsz=6.1)
    txt(sl, 7.42, 4.42, 2.82, 0.13, '子网内为私有端点 NIC · PaaS 本体在 Azure 平台侧', sz=5.9, c=FAINT, align='c')

    seg(sl, 7.14, 3.92, 7.28, 3.92, AZ, 1.4)          # AKS -> PE
    txt(sl, 7.02, 3.70, 0.40, 0.12, '内网', sz=5.6, b=True, c=AZ, align='c')
    txt(sl, 7.02, 4.00, 0.40, 0.12, '调用', sz=5.6, b=True, c=AZ, align='c')

    # users -> WAF
    seg(sl, 2.24, 2.97, 2.52, 2.97, GRN, 1.6)
    txt(sl, 1.30, 2.78, 1.0, 0.13, 'HTTPS 443', sz=6.0, b=True, c=GRN, align='c')
    # source -> AKS extraction
    seg(sl, 6.60, 2.26, 6.60, 3.32, AMB, 1.3, MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 6.74, 3.14, 2.3, 0.13, 'SeaTunnel 抽取 · NAT 固定 IP', sz=6.0, b=True, c=AMB)
    # NAT egress note (arrow up)
    seg(sl, 5.81, 2.72, 5.81, 2.30, AMB, 1.1, MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 4.4, 2.28, 1.35, 0.12, '出站固定 IP ↑', sz=5.9, b=True, c=AMB, align='r')
    # PE -> Databricks private link
    poly(sl, [(10.20, 3.34), (10.20, 2.42), (10.58, 2.42)], PUR, 1.4, MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 8.45, 2.24, 2.0, 0.13, 'Private Link · JDBC / ODBC 443', sz=6.0, b=True, c=PUR, align='r')

    # ---- PaaS band ----
    rect(sl, 0.22, 4.84, 12.90, 1.44, 'F0FAF6', GRN, 1.3, 0.08)
    txt(sl, 0.38, 4.92, 8.0, 0.2, 'Azure PaaS 数据与治理服务 · 全部私有端点接入 · 公网访问已禁用', sz=9, b=True, c=GRN)
    paas = [
        ('Microsoft Entra ID', '身份与单点登录', 'SAML / OIDC · 托管标识', True),
        ('Application Insights', '审计留痕与链路追踪', '操作日志 · 告警规则', False),
        ('Azure Monitor', '资源指标与告警', '仪表板 · 自动化响应', False),
        ('Azure Key Vault', '密钥与凭据托管', 'CMK · 源库连接串', False),
        ('Azure SQL Database', '元数据库 Standard S2', '主备自动故障转移 · PITR', False),
    ]
    xx = 0.34
    for t, role, sub, is_global in paas:
        rect(sl, xx, 5.18, 1.98, 1.00, 'FFFFFF', AZ, 1.1, 0.06, shadow=True)
        rect(sl, xx + 0.02, 5.24, 0.035, 0.88, AZ, rad=0.017)
        txt(sl, xx + 0.12, 5.24, 1.8, 0.16, t, sz=7.4, b=True, c=INK, align='c')
        txt(sl, xx + 0.08, 5.42, 1.84, 0.14, role, sz=6.5, b=True, c=AZ, align='c')
        txt(sl, xx + 0.08, 5.57, 1.84, 0.13, sub, sz=6.0, c=MUT, align='c')
        if is_global:
            chip(sl, xx + 0.34, 5.73, 1.30, 0.17, '全局服务 · 443', None, 'F1F5F9', 'CBD5E1', 0.9, 0.04,
                 tc=SLA, tsz=5.8)
            txt(sl, xx, 5.95, 1.98, 0.14, '认证走 443 · 无私有端点', sz=6.0, b=True, c=MUT, align='c')
        else:
            chip(sl, xx + 0.34, 5.73, 1.30, 0.17, 'Private Endpoint', None, L_BLUE, B_BLUE, 0.9, 0.04,
                 tc=AZD, tsz=5.8)
            txt(sl, xx, 5.95, 1.98, 0.14, '公网访问：已禁用', sz=6.0, b=True, c='DC2626', align='c')
        xx += 2.08
    # storage card (wider)
    rect(sl, 10.74, 5.18, 2.06, 1.00, 'FFFFFF', AZ, 1.1, 0.06, shadow=True)
    rect(sl, 10.76, 5.24, 0.035, 0.88, AZ, rad=0.017)
    txt(sl, 10.80, 5.23, 1.96, 0.15, '两个独立存储账户', sz=7.4, b=True, c=INK, align='c')
    txt(sl, 10.80, 5.39, 1.96, 0.30,
        [dict(runs=[('① ADLS Gen2（HNS）· Iceberg 湖表 · ZRS', {'c': AZ, 'sz': 5.9, 'b': True})], ls=1.15),
         dict(runs=[('② Blob（FNS）· 附件 · Cool→Cold · 可 WORM', {'c': AMB, 'sz': 5.9, 'b': True})], ls=1.15)])
    chip(sl, 11.12, 5.73, 1.30, 0.17, 'Private Endpoint ×2', None, L_BLUE, B_BLUE, 0.9, 0.04, tc=AZD, tsz=5.8)
    txt(sl, 10.74, 5.95, 2.06, 0.14, '公网访问：已禁用', sz=6.0, b=True, c='DC2626', align='c')

    # PE -> PaaS arrows
    seg(sl, 7.57, 4.58, 7.57, 5.16, AZ, 1.2, MSO_LINE_DASH_STYLE.DASH)
    seg(sl, 9.65, 4.58, 9.65, 5.16, AZ, 1.2, MSO_LINE_DASH_STYLE.DASH)
    poly(sl, [(10.10, 4.58), (10.10, 4.90), (11.55, 4.90), (11.55, 5.16)], AZ, 1.2, MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 8.02, 4.88, 1.5, 0.13, 'Private Link 私网流量', sz=6.0, b=True, c=AZ, align='c')
    txt(sl, 8.02, 5.01, 1.5, 0.13, '与上方端点一一对应', sz=6.0, c=MUT, align='c')
    # NCC -> storage
    seg(sl, 12.40, 4.70, 12.40, 5.16, GRN, 2.0)
    txt(sl, 12.48, 4.80, 0.72, 0.34,
        [dict(runs=[('★ NCC 托管', {'c': GRN, 'sz': 5.9, 'b': True})], ls=1.1),
         dict(runs=[('私有端点通道', {'c': GRN, 'sz': 5.9, 'b': True})], ls=1.1)])

    # ---- footer: key points + HA/DR ----
    rect(sl, 0.22, 6.40, 9.50, 0.94, 'FFFFFF', BORDER, 1.1, 0.06)
    txt(sl, 0.36, 6.46, 3.0, 0.17, '关键设计要点', sz=8.5, b=True, c=INK)
    pts_l = [
        '除 Entra ID（全局认证）外，PaaS 全部私有端点接入，流量不出骨干网；入口 WAF + AppGW 收敛',
        'Serverless 不支持 VNet 对等 / VPN —— 私有访问 ADLS 唯一方式：NCC + 托管私有端点',
        '源库抽取走 NAT 固定出口 IP 白名单；跨机房走 ExpressRoute / 站点到站点 VPN',
    ]
    pts_r = [
        '元数据与销毁策略存 SQL DB S2 · 审计进 App Insights · 密钥由 Key Vault 托管',
        '相比方案 B：省去 Trino 集群与双子网注入，查询层完全托管、零集群运维',
        '新建工作区仅 Premium 档可选（Azure 现行政策），UC / PE / NCC 默认可用；NCC 需 Account Admin',
    ]
    for i, s in enumerate(pts_l):
        txt(sl, 0.36, 6.65 + i * 0.155, 4.55, 0.15,
            [dict(runs=[('①②③④⑤⑥'[i] + ' ', {'c': AZ, 'b': True, 'sz': 6.2}), (s, {'c': SUB, 'sz': 6.2})])])
    for i, s in enumerate(pts_r):
        txt(sl, 5.06, 6.65 + i * 0.155, 4.55, 0.15,
            [dict(runs=[('①②③④⑤⑥'[i + 3] + ' ', {'c': AZ, 'b': True, 'sz': 6.2}), (s, {'c': SUB, 'sz': 6.2})])])
    rect(sl, 9.84, 6.40, 3.28, 0.94, L_ROS, 'F87171', 1.0, 0.06, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 9.98, 6.46, 3.0, 0.16, '高可用 & 容灾（HA / DR）', sz=7.5, b=True, c='B91C1C')
    ha = ['AKS 跨 AZ · SQL 主备自动转移 · 存储 ZRS 冗余',
          '平台托管无单点 · Serverless 空闲归零自动重建',
          'SQL 自动备份 + PITR · ADLS 生命周期分层 / GRS',
          'RPO ≤ 15 min · RTO ≤ 2 h · 审计独立备份不销毁']
    for i, s in enumerate(ha):
        txt(sl, 9.98, 6.65 + i * 0.155, 3.02, 0.15,
            [dict(runs=[('· ', {'c': ROS, 'b': True, 'sz': 6.0}), (s, {'c': SUB, 'sz': 6.0})])])

# ================================================== SLIDE 2 · APP ARCH B ====
def slide_app(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 2
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb(BG)
    title_bar(sl, '应用架构图 · 方案 B2（湖仓一体 · Databricks SQL Serverless 读）', '', GRN)

    # band 01
    rect(sl, 0.22, 0.62, 12.90, 0.82, L_GRN, '86EFAC', 1.2, 0.08)
    txt(sl, 0.40, 0.69, 2.0, 0.17, '01 访问入口', sz=8.5, b=True, c=GRND)
    users = [('业务查询员', '查本系统归档数据'), ('审计 / 合规岗', '跨系统检索与举证'), ('平台运维', '监控作业与销毁')]
    xx = 0.40
    for t, s in users:
        chip(sl, xx, 0.92, 2.00, 0.42, t, s, 'FFFFFF', 'BBF7D0', 1.0, 0.05, tsz=8, sc=MUT, ssz=6.4, shadow=True)
        person(sl, xx + 0.22, 1.13, AZ, 0.95)
        xx += 2.16
    chip(sl, 10.44, 0.90, 2.42, 0.46,
         [dict(runs=[('Microsoft Entra ID', {'c': '1D4ED8', 'sz': 8.2, 'b': True})])],
         'SSO / OIDC / 托管标识', L_BLUE, '93C5FD', 1.2, 0.05, sc='1E40AF', ssz=6.5, shadow=True)
    txt(sl, 8.3, 1.02, 2.1, 0.2, '统一身份 · 单点登录 →', sz=6.8, c=MUT, align='r')
    seg(sl, 4.35, 1.44, 4.35, 1.56, CYN, 2.0)
    txt(sl, 4.47, 1.445, 1.2, 0.11, 'HTTPS 443 · SSO', sz=5.8, b=True, c=CYN)

    # band 02
    rect(sl, 0.22, 1.56, 12.90, 1.04, L_CYN, B_CYNB, 1.2, 0.08)
    txt(sl, 0.40, 1.63, 3.0, 0.17, '02 自研应用层（AKS）', sz=8.5, b=True, c=CYND)
    chip(sl, 0.40, 1.86, 2.94, 0.62, 'React 配置化查询平台', 'JSON Schema 渲染查询表单 · 表格 / 明细 / 附件预览 / 导出',
         'FFFFFF', CYN, 1.4, 0.06, tc=INK, tsz=8.2, sc=MUT, ssz=6.4, shadow=True)
    rect(sl, 3.89, 1.86, 5.78, 0.62, 'FFFFFF', CYN, 1.4, 0.06, shadow=True)
    txt(sl, 3.95, 1.86, 5.66, 0.62,
        [dict(runs=[('.NET 10 后端服务组', {'c': INK, 'sz': 8.2, 'b': True})], align='c'),
         dict(runs=[('元数据 · 查询代理 · 权限校验 · 附件代理 · 审计上报 · 销毁调度', {'c': MUT, 'sz': 6.6})],
              align='c', sb=2),
         dict(runs=[('后端二次鉴权 · 参数化 SQL · 私有端点访问下游', {'c': FAINT, 'sz': 6.6})], align='c', sb=1)],
        anchor='m')
    seg(sl, 3.40, 2.17, 3.83, 2.17, CYN, 1.8)   # React -> .NET
    chip(sl, 10.22, 1.86, 2.64, 0.62, 'Azure SQL 元数据库', '系统台账 / 表清单 / 保留期 / 销毁策略 / 审计索引',
         'FFFFFF', CYN, 1.4, 0.06, tc=INK, tsz=8.2, sc=MUT, ssz=6.4, shadow=True)
    seg(sl, 9.73, 2.17, 10.16, 2.17, CYN, 1.8)
    txt(sl, 9.14, 2.00, 1.6, 0.14, '读写配置', sz=6.4, b=True, c=CYN, align='c')

    # band 03
    rect(sl, 0.22, 2.74, 12.90, 2.42, 'FFFFFF', BORDER, 1.2, 0.08)
    txt(sl, 0.40, 2.81, 3.4, 0.17, '03 核心链路（写入 / 查询 / 附件）', sz=8.5, b=True, c='334155')
    lanes = [
        (0.40, 'WRITE PATH', GRND, L_GRNB, B_GRNB, '归档写入 / 到期销毁',
         [('Azure Data Factory', '编排调度', 'FDBA74'), ('SeaTunnel', '结构化抽取', 'FDBA74')],
         'DCFCE7', B_GRNB, GRND, 'Databricks Serverless：ETL · Compaction · DROP PARTITION + VACUUM'),
        (4.62, 'READ PATH', PURD, L_PURB, B_PURB, '在线查询 / 跨系统检索',
         [('查询代理', '限流 / 超时 / 审计', 'C4B5FD'), ('Databricks SQL Serverless', '读侧查询 · 空闲缩 0', 'A78BFA')],
         'EDE9FE', B_PURB, PURD, 'Unity Catalog 表级授权 · SERVE 层首屏 2–5 s'),
        (8.84, 'ATTACHMENT PATH', 'B45309', L_AMB, B_AMB, '非结构化附件 / 预览下载',
         [('Storage Mover', '迁移附件', 'FCD34D'), ('附件代理', '鉴权 / 流式回吐', 'FCD34D')],
         'FEF3C7', B_AMB, '92400E', 'SHA-256 台账 · 附件索引表 · 后端审计'),
    ]
    for lx, tag, tagc, lf, bl, ttl, boxes, pf, pb, ptc, ptext in lanes:
        rect(sl, lx, 3.02, 4.02, 1.94, lf, bl, 1.2, 0.08)
        txt(sl, lx + 0.16, 3.10, 2.4, 0.14, tag, sz=6.4, b=True, c=tagc, spc=140)
        txt(sl, lx + 0.16, 3.26, 3.7, 0.2, ttl, sz=9.2, b=True, c=INK)
        bw = 1.72
        chip(sl, lx + 0.14, 3.56, bw, 0.55, boxes[0][0], boxes[0][1], 'FFFFFF', boxes[0][2], 1.2, 0.05,
             tsz=7.8, sc=MUT, ssz=6.3, shadow=True)
        chip(sl, lx + 0.16 + bw + 0.30, 3.56, bw, 0.55, boxes[1][0], boxes[1][1], 'FFFFFF', boxes[1][2], 1.2, 0.05,
             tsz=7.8, sc=MUT, ssz=6.3, shadow=True)
        ac = {'WRITE PATH': AMB, 'READ PATH': PUR, 'ATTACHMENT PATH': AMB}[tag]
        seg(sl, lx + 0.14 + bw + 0.04, 3.835, lx + 0.14 + bw + 0.28, 3.835, ac, 2.0)
        rect(sl, lx + 0.16, 4.28, 3.70, 0.52, pf, pb, 1.1, 0.06)
        txt(sl, lx + 0.22, 4.28, 3.58, 0.52, ptext, sz=6.8, b=True, c=ptc, align='c', anchor='m')

    # app -> lanes arrows
    poly(sl, [(4.20, 2.48), (4.20, 2.90), (2.41, 2.90), (2.41, 3.00)], GRN, 2.2)
    txt(sl, 2.62, 2.73, 1.9, 0.14, '接入配置 / 销毁调度', sz=6.4, b=True, c=GRN)
    seg(sl, 6.63, 2.48, 6.63, 3.00, PUR, 2.2)
    txt(sl, 6.74, 2.66, 1.3, 0.14, '查询请求', sz=6.4, b=True, c=PUR)
    poly(sl, [(7.90, 2.48), (7.90, 2.90), (10.85, 2.90), (10.85, 3.00)], AMB, 2.0)
    txt(sl, 8.9, 2.73, 1.8, 0.14, '附件预览 / 下载', sz=6.4, b=True, c=AMB)

    # band 04
    rect(sl, 0.22, 5.28, 12.90, 0.82, L_BLUE, B_BLUE, 1.2, 0.08)
    txt(sl, 0.40, 5.35, 4.0, 0.17, '04 数据底座与统一治理', sz=8.5, b=True, c='1D4ED8')
    chip(sl, 0.40, 5.56, 3.90, 0.46, 'ADLS Gen2 · Iceberg 湖表', 'RAW / CURATED / LAKE / SERVE · 1:1 原样归档',
         'FFFFFF', '60A5FA', 1.2, 0.05, tc=INK, tsz=8, sc=MUT, ssz=6.3, shadow=True)
    chip(sl, 4.50, 5.56, 3.60, 0.46, 'Unity Catalog', '表级权限 / 血缘 / 快照回溯',
         'FFFFFF', '60A5FA', 1.2, 0.05, tc=INK, tsz=8, sc=MUT, ssz=6.3, shadow=True)
    chip(sl, 8.30, 5.56, 4.56, 0.46, 'Azure Blob · 附件账户', 'Cool / Cold 分层 · WORM 可选 · 与业务表分账户',
         'FFFFFF', 'F59E0B', 1.2, 0.05, tc=INK, tsz=8, sc=MUT, ssz=6.3, shadow=True)
    # lanes -> data arrows
    seg(sl, 2.41, 4.96, 2.41, 5.54, GRN, 2.2)
    txt(sl, 2.52, 5.175, 1.4, 0.11, '写入 Iceberg', sz=6.4, b=True, c=GRN)
    poly(sl, [(6.63, 4.96), (6.63, 5.22), (3.60, 5.22), (3.60, 5.54)], PUR, 2.2)
    txt(sl, 4.62, 5.00, 1.8, 0.14, '读 Iceberg / ABFS', sz=6.4, b=True, c=PUR)
    seg(sl, 10.85, 4.96, 10.85, 5.54, AMB, 2.0)
    txt(sl, 10.96, 5.175, 1.3, 0.11, '附件取流', sz=6.4, b=True, c=AMB)

    # footer pills
    rect(sl, 0.40, 6.24, 8.02, 0.40, 'FFFFFF', 'CBD5E1', 1.0, 0.06)
    txt(sl, 0.56, 6.24, 7.7, 0.40,
        [dict(runs=[('安全底座：', {'b': True, 'c': '334155', 'sz': 6.8}),
                    ('Key Vault · App Insights · Azure Monitor · Databricks Query History · 全部 PaaS 私有端点 · 入口 WAF',
                     {'c': SUB, 'sz': 6.8})])], anchor='m')
    rect(sl, 8.64, 6.24, 4.22, 0.40, 'FFF7ED', 'FDBA74', 1.0, 0.06)
    txt(sl, 8.82, 6.24, 3.86, 0.40,
        [dict(runs=[('B2 读侧：', {'b': True, 'c': 'B45309', 'sz': 6.8}),
                    ('Databricks SQL Serverless 完全托管 · 空闲缩 0；代价是引擎绑定（数据仍为开放 Iceberg 可迁移）',
                     {'c': AMBD, 'sz': 6.8})])], anchor='m')
    # legend
    leg = [('自研应用组件（AKS）', CYN), ('写侧数据链路', GRN), ('读侧查询链路', PUR),
           ('编排 / 抽取', AMB), ('Azure PaaS 数据服务', AZ)]
    xx = 0.40
    for t, c in leg:
        oval(sl, xx, 6.83, 0.09, 0.09, c)
        txt(sl, xx + 0.14, 6.78, 1.9, 0.16, t, sz=6.4, c=MUT)
        xx += 0.32 + (len(t) * 0.082 + 0.5)

# ============================================ SLIDE 3 · LAKEHOUSE (存算分离) =
def vstage(sl, y, h, sid, col, t, en, subs, fill, bl, lw=1.2):
    """竖向流水线的单个阶段卡：左侧脊柱圆点 + 标题块，右侧内容区（x2.44–8.45）。"""
    rect(sl, 0.78, y, 7.77, h, fill, bl, lw, 0.08)
    oval(sl, 0.38, y + h / 2 - 0.17, 0.34, 0.34, 'FFFFFF', col, 1.8)
    txt(sl, 0.38, y + h / 2 - 0.17, 0.34, 0.34, sid, sz=6.5 if len(sid) <= 2 else 5.3,
        b=True, c=col, align='c', anchor='m', mono=True)
    txt(sl, 0.90, y + 0.06, 1.46, 0.18, t, sz=9.2, b=True, c=INK)
    txt(sl, 0.90, y + 0.25, 1.46, 0.12, en, sz=5.2, b=True, c=col, spc=110, mono=True)
    yy = y + 0.40
    for s, strong, cc in subs:
        txt(sl, 0.90, yy, 1.48, 0.125, s, sz=5.9, c=cc or MUT, b=strong)
        yy += 0.13

def slide_lake(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 3
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb(BG)
    title_bar(sl, '方案 B2 · 湖仓一体 · 存算分离架构', '', PUR)

    # ================= 左侧：竖向主管道 =================
    seg(sl, 0.55, 0.74, 0.55, 6.90, 'E2E8F0', 2.5, arrow=False)   # 脊柱

    # ---- 01 源系统 ----
    vstage(sl, 0.62, 0.68, '01', SLA, '源系统', 'SOURCE',
           [('退役 · 只读', False, None), ('15 套待归档', False, None)], 'F6F8FA', 'CBD5E1')
    chip(sl, 2.44, 0.70, 2.94, 0.52, '结构化数据',
         'Oracle · MySQL · DB2 / 其他 · JDBC 读取 · 约 100 张表',
         'FFFFFF', SLA, 1.3, 0.05, tc=INK, tsz=7.4, sc=MUT, ssz=6.0)
    chip(sl, 5.51, 0.70, 2.94, 0.52, '非结构化附件',
         'NFS · S3 · Azure Blob · PDF / 图片 / 影像件',
         L_AMB, AMB, 1.3, 0.05, tc=AMBD, tsz=7.4, sc='78350F', ssz=6.0)
    seg(sl, 5.45, 1.30, 5.45, 1.41, '94A3B8', 1.6)

    # ---- B-00 接入配置 ----
    vstage(sl, 1.42, 0.92, 'B-00', CYN, '接入配置', 'CONFIGURE',
           [('页面操作 · 零代码', False, None), ('配置存元数据库', False, None),
            ('新系统接入复用', True, CYN)], L_CYN, B_CYNB)
    b00 = [
        ('① 录入连接信息', ['主机 / 端口 / 账号 · 凭据存 Key Vault', '「测试连接」校验连通与权限']),
        ('② 自动读取表清单', ['扫描 information_schema · 表名 / 行数 / 大小', '共发现 10 张表（含类型与主键）']),
        ('③ 勾选待迁移表', ['☑ 订单主表 · 明细 · 客户 · 商品 · 支付', '☐ 临时 / 日志 / 中间 / 缓存 / 备份']),
    ]
    bx = 2.44
    for t, lines in b00:
        rect(sl, bx, 1.48, 1.85, 0.64, 'FFFFFF', CYN, 1.3, 0.05)
        txt(sl, bx + 0.10, 1.54, 1.68, 0.13, t, sz=6.8, b=True, c=INK)
        yy = 1.70
        for l in lines:
            txt(sl, bx + 0.10, yy, 1.68, 0.22, l, sz=5.7, c=MUT, leading=1.1)
            yy += 0.21
        bx += 2.08
    seg(sl, 4.30, 1.80, 4.50, 1.80, CYN, 1.5)
    seg(sl, 6.38, 1.80, 6.58, 1.80, CYN, 1.5)
    txt(sl, 2.44, 2.15, 6.01, 0.13, '按配置生成 SeaTunnel 任务 · 仅入湖勾选的 5 张表（10 选 5）',
        sz=6.2, b=True, c=CYN, align='c')
    seg(sl, 5.45, 2.34, 5.45, 2.45, '94A3B8', 1.6)

    # ---- 02 数据抽取 ----
    vstage(sl, 2.46, 0.68, '02', ORG, '数据抽取', 'INGEST',
           [('批量 · 全量 + 增量', False, None), ('ADF 编排调度', False, None)], 'FFF7ED', 'FED7AA')
    chip(sl, 2.44, 2.54, 2.94, 0.52, '结构化 · SeaTunnel',
         '200+ 连接器 · JDBC · 类型映射 · 对账 → Iceberg 表',
         'FFFFFF', ORG, 1.3, 0.05, tc=INK, tsz=7.4, sc=MUT, ssz=6.0)
    chip(sl, 5.51, 2.54, 2.94, 0.52, '附件 · Storage Mover',
         'PDF / 图片 / Office 影像 · 服务免费 → ADLS 附件容器',
         L_AMB, AMB, 1.5, 0.05, tc=AMBD, tsz=7.4, sc='78350F', ssz=6.0)
    seg(sl, 5.45, 3.14, 5.45, 3.25, '94A3B8', 1.6)

    # ---- B-03 Iceberg 冷湖 ----
    vstage(sl, 3.26, 1.24, 'B-03', PUR, 'Iceberg 冷湖', 'COLD LAKE · STORAGE',
           [('ADLS Gen2 · Parquet + ZSTD', False, None), ('≈¥90–150/TB·月 · PB 级', True, PUR),
            ('100 张表 · 1:1 原样', False, None), ('Hot / Cool 分层 · ZRS', False, None)], 'FAF9FF', B_PUR)
    txt(sl, 2.44, 3.32, 6.01, 0.14, 'Iceberg 表 · 四层分层（RAW → SERVE 逐层升温，越往下越热）',
        sz=6.4, b=True, c=PUR)
    layers = [
        ('RAW · 原始落地（最冷）', 'SeaTunnel 全量 1:1 保真 · 分区 sys_id / ingest_date', 'EEF2FF', 'C7D2FE', '818CF8', '4338CA'),
        ('CURATED · 清洗规整', '类型统一 · 缺失值处理 · SHA-256 台账 · 分区归约', 'EDE9FE', 'C4B5FD', 'A78BFA', '6D28D9'),
        ('LAKE · 建模归档', '业务维度建模 · 生命周期分区（销毁友好）', 'F5F3FF', 'A78BFA', '7C3AED', '5B21B6'),
        ('SERVE · 加速服务（最热）', '热表物化 · Z-Order · 供 Databricks SQL Serverless 直查', 'FEF3C7', 'FCD34D', 'D97706', 'B45309'),
    ]
    ly = 3.50
    for t, s, f, bl, ac, tc in layers:
        rect(sl, 2.44, ly, 6.01, 0.225, f, bl, 1.0, 0.04)
        rect(sl, 2.46, ly + 0.025, 0.035, 0.175, ac, rad=0.017)
        txt(sl, 2.58, ly, 5.80, 0.225,
            [dict(runs=[(t, {'c': tc, 'sz': 7.0, 'b': True}),
                        ('　' + s, {'c': MUT, 'sz': 6.0})])], anchor='m')
        ly += 0.245
    # 存算分离边界
    seg(sl, 0.78, 4.65, 8.55, 4.65, '94A3B8', 1.6, MSO_LINE_DASH_STYLE.DASH, arrow=False)
    rect(sl, 4.62, 4.55, 1.76, 0.20, 'FFFFFF', 'CBD5E1', 1.0, 0.06)
    txt(sl, 4.62, 4.55, 1.76, 0.20, '━━ 存算分离边界 ━━', sz=6.8, b=True, c=SLA, align='c', anchor='m')
    seg(sl, 3.40, 4.50, 3.40, 4.80, PUR, 1.8)          # 湖 -> Serverless 读
    txt(sl, 3.52, 4.51, 1.00, 0.12, 'Serverless 直查', sz=5.6, b=True, c=PUR)
    seg(sl, 7.80, 4.80, 7.80, 4.50, GRN, 1.8)          # Databricks ETL 写回
    txt(sl, 6.45, 4.67, 1.28, 0.12, 'ETL 写入 · Databricks', sz=5.6, b=True, c=GRN, align='r')

    # ---- B-04 统一计算层 ----
    vstage(sl, 4.80, 0.84, 'B-04', GRN, '统一计算层', 'COMPUTE',
           [('Databricks 读写一体', False, None), ('弹性 · 按量计费', False, None)], L_GRNB, 'A7F3D0')
    comp = [
        (2.44, 'Databricks · 写侧 ETL',
         ['RAW → CURATED → LAKE → SERVE', 'Spark 批处理 · Iceberg 写入 · Z-Order',
          'Compaction · VACUUM · 到期销毁', '按 DBU 计费 · 按需拉起 ≈¥0.4–0.8 万/月']),
        (5.51, 'Databricks SQL Serverless',
         ['读侧查询 · 零运维 · 故障自动重建', 'Iceberg 直查 · SERVE 加速 · 首屏 2–5 s',
          'UC 表 / 行 / 列级授权 + 审计', 'X-Small 6 DBU/h · 空闲缩 0 ≈¥0.5–0.9 万/月']),
    ]
    for cx0, t, lines in comp:
        rect(sl, cx0, 4.86, 2.94, 0.72, 'FFFFFF', GRN, 1.3, 0.05)
        txt(sl, cx0 + 0.10, 4.92, 2.76, 0.13, t, sz=7.2, b=True, c=INK, align='c')
        yy = 5.08
        for i, l in enumerate(lines):
            hot = (i == 3)
            txt(sl, cx0 + 0.08, yy, 2.80, 0.115, l, sz=5.9, c=GRN if hot else MUT, b=hot, align='c')
            yy += 0.125
    seg(sl, 5.45, 5.64, 5.45, 5.75, '94A3B8', 1.6)

    # ---- B-05 消费入口 ----
    vstage(sl, 5.76, 0.72, 'B-05', CYN, '消费入口', 'CONSUME',
           [('元数据驱动', False, None), ('零代码接入', False, None)], L_CYN, 'A5F3FC')
    rect(sl, 2.44, 5.82, 6.01, 0.60, 'FFFFFF', CYN, 1.4, 0.05, shadow=True)
    txt(sl, 2.56, 5.88, 5.77, 0.50,
        [dict(runs=[('元数据驱动 · 动态配置化查询平台', {'c': INK, 'sz': 7.6, 'b': True})], align='c'),
         dict(runs=[('JSON 配置定义查询 / 筛选 / 列 / 权限 · 前端动态渲染 · 一次研发处处复用',
                     {'c': MUT, 'sz': 6.0})], align='c', sb=3),
         dict(runs=[('新系统接入零代码 · 边际 +3~6 人天 / 系统 · 首期 60 人天 · ≥15 系统适用',
                     {'c': CYN, 'sz': 6.0, 'b': True})], align='c', sb=2)])

    # ---- Retention 销毁回路 ----
    rect(sl, 0.78, 6.58, 7.77, 0.72, L_ROS, ROS, 1.2, 0.08, dash=MSO_LINE_DASH_STYLE.DASH)
    oval(sl, 0.38, 6.77, 0.34, 0.34, 'FFFFFF', ROS, 1.8)
    txt(sl, 0.38, 6.77, 0.34, 0.34, '↺', sz=10, b=True, c=ROS, align='c', anchor='m')
    txt(sl, 0.90, 6.66, 1.46, 0.16, 'Retention', sz=8.4, b=True, c=ROS)
    txt(sl, 0.90, 6.84, 1.46, 0.11, 'RETENTION', sz=5.2, b=True, c=ROS, spc=110, mono=True)
    txt(sl, 0.90, 6.99, 1.48, 0.24, '作用于 B-03 分区\n回写台账 · 附件同步删', sz=5.8, c=MUT, leading=1.15)
    rsteps = [
        ('① 到期扫描', '扫描 Iceberg meta'), ('② DROP PARTITION', '秒级 · Iceberg 原生'),
        ('③ VACUUM', '物理清除快照'), ('④ 附件 + 审计', '同步删除 · 留痕'),
    ]
    rx = 2.44
    for t, s in rsteps:
        chip(sl, rx, 6.68, 1.34, 0.42, t, s, 'FFFFFF', ROS, 1.2, 0.05, tc=INK, tsz=6.6, sc=MUT, ssz=5.6)
        rx += 1.56
    for ax in (3.81, 5.37, 6.93):
        seg(sl, ax, 6.89, ax + 0.16, 6.89, ROS, 1.5)
    txt(sl, 2.44, 7.14, 6.01, 0.13, '审计与元数据独立备份 · 不随业务数据销毁 · 纳入年度销毁演练',
        sz=5.8, b=True, c='9F1239', align='c')
    # 回路箭头：从 Retention 左侧绕到 B-03 左缘
    poly(sl, [(0.78, 7.12), (0.30, 7.12), (0.30, 4.20), (0.76, 4.20)], ROS, 1.2, MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 0.10, 4.85, 0.20, 1.30, '作用于冷湖分区', sz=5.6, b=True, c=ROS, align='c', anchor='m', vert=True)

    # ================= 右侧：附件轨道 / 成本 / 要点 =================
    # R1 附件对象存储
    rect(sl, 8.75, 0.62, 4.36, 1.58, L_AMB, AMB, 1.1, 0.07, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 8.89, 0.70, 4.1, 0.16, [dict(runs=[
        ('附件对象存储 · ', {'c': AMB, 'sz': 8.2, 'b': True}),
        ('独立存储账户（与湖分开）', {'c': AMBD, 'sz': 8.2, 'b': True})])])
    txt(sl, 8.89, 0.88, 4.1, 0.13, '贯穿 02 抽取 → B-03 存储 · 供 B-05 直接消费', sz=5.9, c='A16207')
    r1chips = [('FNS 平面命名空间', '非 HNS · 按次计费（不按 4MB 切块）'), ('Cool → Cold 分层', '生命周期自动降档'),
               ('SHA-256 台账', '前缀 sysXX/yyyy-mm/'), ('后端签 SAS', '前端直连下载')]
    for i, (t, s) in enumerate(r1chips):
        cx = 8.87 + (i % 2) * 2.12; cy = 1.06 + (i // 2) * 0.50
        chip(sl, cx, cy, 2.05, 0.44, t, s, 'FFFFFF', 'FDE68A', 1.0, 0.05, tc=INK, tsz=6.8, sc='A16207', ssz=5.8)
    seg(sl, 8.75, 2.80, 8.57, 2.80, AMB, 1.3, MSO_LINE_DASH_STYLE.DASH)      # 入池
    seg(sl, 8.75, 3.88, 8.57, 3.88, AMB, 1.3, MSO_LINE_DASH_STYLE.DASH)      # 台账共池
    poly(sl, [(10.93, 2.20), (10.93, 2.34), (8.66, 2.34), (8.66, 6.12), (8.57, 6.12)],
         AMB, 1.3, MSO_LINE_DASH_STYLE.DASH)                                  # URL 直连消费
    txt(sl, 8.82, 2.38, 2.1, 0.12, '附件 URL 直连消费（SAS）', sz=5.9, b=True, c=AMB)

    # R2 存算分离 · 成本结构
    rect(sl, 8.75, 2.56, 4.36, 1.34, 'FFFFFF', B_PUR, 1.2, 0.07)
    txt(sl, 8.89, 2.64, 4.1, 0.16, '存算分离 · 成本结构', sz=8.2, b=True, c=PURD)
    chip(sl, 8.89, 2.88, 1.42, 0.34, '存储侧 · 常驻', None, 'F5F3FF', B_PUR, 1.0, 0.05, tc=PURD, tsz=6.6)
    txt(sl, 10.41, 2.88, 2.62, 0.36,
        [dict(runs=[('ADLS Hot/Cool ≈¥90–150/TB·月 · PB 级从容 · ZRS', {'c': SUB, 'sz': 5.9})], ls=1.15)])
    chip(sl, 8.89, 3.30, 1.42, 0.34, '计算侧 · 按需', None, 'F0FDF4', 'A7F3D0', 1.0, 0.05, tc=GRND, tsz=6.6)
    txt(sl, 10.41, 3.30, 2.62, 0.36,
        [dict(runs=[('写侧 ETL ≈¥0.4–0.8 万/月 · 读侧查询 ≈¥0.5–0.9 万/月', {'c': SUB, 'sz': 5.9})], ls=1.15)])
    txt(sl, 8.89, 3.72, 4.1, 0.13, '存储与计算独立扩缩 · 空闲缩到 0 · 只为使用的算力付费',
        sz=5.9, b=True, c=PUR, align='c')

    # R3 方案 B 要点
    rect(sl, 8.75, 4.02, 4.36, 1.72, 'FFFFFF', BORDER, 1.2, 0.07)
    txt(sl, 8.89, 4.10, 4.1, 0.16, '方案 B2 要点', sz=8.2, b=True, c=INK)
    blts = [
        ('存算分离', '冷湖常驻，Serverless 读写按需拉起 · 空闲缩 0', PUR),
        ('到期即毁', 'DROP PARTITION + VACUUM 物理清除，附件与审计同步', ROS),
        ('零代码接入', '元数据驱动前端 · 边际 +3~6 人天 / 系统', CYN),
        ('读侧托管', 'Databricks SQL Serverless · 零集群运维 · UC 授权', GRN),
        ('适用规模', '≥15 套系统或 >20 TB 长期保留 · PB 级从容', SLA),
    ]
    yy = 4.32
    for tag, s, c in blts:
        txt(sl, 8.89, yy, 4.12, 0.24,
            [dict(runs=[('· ' + tag + ' — ', {'c': c, 'sz': 6.1, 'b': True}),
                        (s, {'c': SUB, 'sz': 6.1})], ls=1.1)])
        yy += 0.265

    # R4 B 的取舍
    rect(sl, 8.75, 5.86, 4.36, 0.86, 'FFF7ED', 'FDBA74', 1.1, 0.07)
    txt(sl, 8.89, 5.94, 4.1, 0.15, 'B2 的取舍', sz=8.0, b=True, c='B45309')
    txt(sl, 8.89, 6.12, 4.12, 0.52,
        [dict(runs=[('读侧完全托管、UC 治理原生；代价是计算引擎绑定 Databricks（数据仍为开放 Iceberg，随时可换引擎）。Premium 为新建工作区唯一档（Azure 现行政策）。',
                     {'c': AMBD, 'sz': 6.1})], ls=1.25)])


# ===================================================== SLIDE 5 · ENCRYPTION =
def col_head(sl, x, y, w, t):
    rect(sl, x, y, w, 0.30, 'FFFFFF', BORDER, 1.1, 0.06)
    txt(sl, x + 0.14, y, w - 0.28, 0.30, t, sz=8.5, b=True, c=INK, anchor='m')

def slide_enc(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 4
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb(BG)
    title_bar(sl, '方案 B2 · 数据加密设计', '', AZ)

    # ================= 左：四层加密模型 + 密钥管理 =================
    col_head(sl, 0.22, 0.62, 4.30, '四层加密模型 + 密钥管理（L1 → L5）')
    layers = [
        ('L1', AZ, L_BLUE, B_BLUE, '存储层', '防磁盘 / 桶泄露',
         'ADLS 账户级 SSE（AES-256）+ Key Vault CMK；数据 / 元数据 / 附件一套密钥策略，可选基础设施双重加密'),
        ('L2', GRN, L_GRN, B_GRN, '传输层', '防链路窃听',
         '全链路 TLS 1.2+：HTTPS 443 入口 · 私有端点私网 · JDBC / ODBC；流量不出 Azure 骨干网'),
        ('L3', ROS, L_ROS, B_ROS, '字段层', '防 DBA / 运维内鬼 ★ 合规核心',
         'aes_encrypt() 落盘即密文 + UC 列掩码查询遮蔽——两者必须组合；GCM 每次随机 IV'),
        ('L4', PUR, L_PUR, B_PUR, '展示层', '防越权查看',
         'UC COLUMN MASK / ROW FILTER 绑定基表，任何访问路径强制生效；B2 原生，Trino 读完全不可用'),
        ('L5', AMB, L_AMB, B_AMB, '密钥层', '防泄露无法止损',
         'Key Vault CMK 自动轮转；信封加密——轮主密钥只重包装 KEK、零数据重写；密文自带 key_version'),
    ]
    yy = 1.00
    for tag, col, lf, bl, name, threat, d in layers:
        rect(sl, 0.22, yy, 4.30, 0.82, lf, bl, 1.6 if tag == 'L3' else 1.1, 0.06)
        rect(sl, 0.36, yy + 0.13, 0.52, 0.28, 'FFFFFF', bl, 1.1, 0.05)
        txt(sl, 0.36, yy + 0.13, 0.52, 0.28, tag, sz=7.5, b=True, c=col, align='c', anchor='m', mono=True)
        txt(sl, 1.02, yy + 0.09, 3.40, 0.15,
            [dict(runs=[(name, {'b': True, 'sz': 7.9, 'c': col if tag == 'L3' else INK}),
                        (' · ' + threat, {'sz': 5.9, 'c': MUT})])])
        txt(sl, 1.02, yy + 0.28, 3.38, 0.48, d, sz=5.9, c=SUB, leading=1.22)
        yy += 0.90

    # ================= 中：字段级加密 · 推荐路线四步 =================
    col_head(sl, 4.66, 0.62, 4.50, '字段级加密 · 推荐路线：aes_encrypt + UC 列掩码')
    steps = [
        ('STEP 1', '密钥入 Secret Scope', 'Key Vault 撑腰 · 托管标识免密拉取，密钥不进代码 / 日志'),
        ('STEP 2', 'ETL 写侧加密', dict(runs=[('aes_encrypt(col, secret(...))', {'mono': True, 'c': CYND, 'sz': 5.9}),
                                               (' · 默认 GCM 随机 IV，绝不能用 ECB', {'c': MUT, 'sz': 5.9})])),
        ('STEP 3', '定义掩码函数', dict(runs=[('is_account_group_member()', {'mono': True, 'c': CYND, 'sz': 5.9}),
                                               (' ? try_aes_decrypt : ***', {'mono': True, 'c': MUT, 'sz': 5.9})])),
        ('STEP 4', '挂载到列', dict(runs=[('ALTER TABLE … SET MASK', {'mono': True, 'c': CYND, 'sz': 5.9}),
                                               ('（必须打在基表，不支持视图）', {'c': MUT, 'sz': 5.9})])),
    ]
    yy = 0.98
    for i, (tag, t, d) in enumerate(steps):
        rect(sl, 4.66, yy, 4.50, 0.60, 'FFFFFF', B_GRN, 1.2, 0.06, shadow=True)
        oval(sl, 4.78, yy + 0.15, 0.30, 0.30, 'FFFFFF', GRN, 1.6)
        txt(sl, 4.78, yy + 0.15, 0.30, 0.30, str(i + 1), sz=8, b=True, c=GRND, align='c', anchor='m', mono=True)
        txt(sl, 5.20, yy + 0.08, 3.86, 0.15, t, sz=7.6, b=True, c=INK)
        txt(sl, 5.20, yy + 0.28, 3.86, 0.26, [d] if isinstance(d, dict) else [(d, {'c': MUT, 'sz': 5.9})],
            leading=1.2)
        if yy < 2.9:
            seg(sl, 6.90, yy + 0.61, 6.90, yy + 0.69, GRN, 1.4)
        yy += 0.70
    rect(sl, 4.66, 3.80, 4.50, 0.40, 'DCFCE7', B_GRN, 1.1, 0.06)
    txt(sl, 4.76, 3.80, 4.30, 0.40, '同一条 SQL：授权者见明文 · 其他人见 *** · 直读 Parquet 只见密文',
        sz=6.8, b=True, c=GRND, align='c', anchor='m')
    rect(sl, 4.66, 4.26, 4.50, 0.28, 'F8FAFC', BORDER, 1.0, 0.05)
    txt(sl, 4.78, 4.26, 4.28, 0.28,
        [dict(runs=[('为何安全：', {'b': True, 'c': INK, 'sz': 5.9}),
                    ('掩码函数以定义者权限取密钥，分析师自己写 aes_decrypt 也拿不到 key',
                     {'c': MUT, 'sz': 5.9})])], anchor='m')
    rect(sl, 4.66, 4.62, 4.50, 0.80, 'FFF7ED', 'FDBA74', 1.1, 0.06)
    txt(sl, 4.78, 4.68, 4.3, 0.14, '三个必守约束', sz=7.2, b=True, c='B45309')
    txt(sl, 4.78, 4.86, 4.28, 0.52,
        [dict(runs=[('❶ 绝不能改成 ECB（暴露频率分布）　❷ 业务用户对 Secret Scope 必须无 READ 权限，否则防线归零',
                     {'c': AMBD, 'sz': 5.8})], ls=1.25),
         dict(runs=[('❸ 掩码不支持视图须挂基表；卸载顺序：先 DROP MASK 再 DROP FUNCTION', {'c': AMBD, 'sz': 5.8})], ls=1.25)])

    # ================= 右：密钥体系 · 合规与销毁 =================
    col_head(sl, 9.32, 0.62, 3.80, '密钥体系 · 合规与销毁')
    rect(sl, 9.32, 0.98, 3.80, 1.72, 'FFFFFF', AZ, 1.2, 0.06)
    chip(sl, 9.46, 1.08, 3.52, 0.34, 'Key Vault · CMK / KEK 主密钥', None, L_BLUE, B_BLUE, 1.1, 0.05,
         tc=AZD, tsz=6.8, shadow=True)
    seg(sl, 11.22, 1.43, 11.22, 1.55, AZ, 1.3)
    txt(sl, 11.32, 1.42, 1.0, 0.12, '包装 / 解包', sz=5.4, c=AZ, b=True)
    chip(sl, 9.46, 1.56, 3.52, 0.34, 'DEK 数据密钥（带 key_version）', None, L_PUR, B_PUR, 1.1, 0.05,
         tc=PURD, tsz=6.8, shadow=True)
    seg(sl, 11.22, 1.91, 11.22, 2.03, PUR, 1.3)
    txt(sl, 11.32, 1.90, 1.0, 0.12, '加密 / 解密', sz=5.4, c=PUR, b=True)
    chip(sl, 9.46, 2.04, 3.52, 0.34, 'Iceberg 密文 + UC 列掩码', None, L_GRN, B_GRN, 1.1, 0.05,
         tc=GRND, tsz=6.8, shadow=True)
    txt(sl, 9.42, 2.44, 3.60, 0.22, '信封加密：轮转主密钥只重新包装 DEK，数据文件一个字节不动',
        sz=5.8, b=True, c=AZ, align='c')
    rect(sl, 9.32, 2.80, 3.80, 0.94, L_ROS, ROS, 1.2, 0.06, dash=MSO_LINE_DASH_STYLE.DASH)
    txt(sl, 9.46, 2.87, 3.5, 0.15, 'Crypto-Shredding · 被遗忘权', sz=7.6, b=True, c=ROS)
    txt(sl, 9.46, 3.06, 3.54, 0.60,
        [dict(runs=[('按数据主体分配独立 DEK——「删除」= 销毁该 DEK，密文永久不可读，无需重写 PB 级文件',
                     {'c': SUB, 'sz': 5.9})], ls=1.25),
         dict(runs=[('归档场景性价比最高的合规设计，建议第一版纳入', {'c': ROS, 'sz': 5.9, 'b': True})], ls=1.25)])
    rect(sl, 9.32, 3.84, 3.80, 1.58, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, 9.46, 3.91, 3.5, 0.15, '合规映射', sz=7.6, b=True, c=INK)
    comp = [
        ('等保 2.0 三级', L_GRN, B_GRN, GRND, '存储 / 传输保密达标；策略在引擎内、审计落表，举证材料天然齐备'),
        ('个保法第 51 条', L_AMB, B_AMB, 'B45309', 'SSE ≠ 去标识化——必须 L3 字段级加密 + L4 强制脱敏，只做 L1 不满足'),
        ('GDPR 被遗忘权', L_ROS, B_ROS, ROS, 'Crypto-Shredding：销毁主体 DEK 即永久不可读'),
    ]
    yy = 4.12
    for tag, lf, bl, tc, d in comp:
        chip(sl, 9.46, yy, 1.06, 0.24, tag, None, lf, bl, 1.0, 0.04, tc=tc, tsz=5.9)
        txt(sl, 10.62, yy - 0.01, 2.42, 0.42, d, sz=5.7, c=MUT, leading=1.2)
        yy += 0.44

    # ================= 底部：三个结构性优势 + 结论 =================
    adv = [
        ('① 授权强制力', '059669', L_BLUE, B_BLUE, '掩码 / 行过滤绑定引擎元数据，任何客户端绕不过；应用层脱敏绕过 API 即失效'),
        ('② 轮转成本', '059669', L_GRN, B_GRN, '信封加密只重包装 KEK、零数据重写；应用层方案轮转需重刷数十 TB 存量'),
        ('③ 实施成本', '7C3AED', L_PUR, B_PUR, '约 2~3 人天（应用层方案 8~12 人天）· 加密字段仍可检索 · 审计平台原生'),
    ]
    xx = 0.22
    for t, chex, lf, bl, d in adv:
        rect(sl, xx, 5.56, 4.20, 0.98, lf, bl, 1.2, 0.06, shadow=True)
        txt(sl, xx + 0.16, 5.68, 3.9, 0.15, t, sz=7.8, b=True, c=chex)
        txt(sl, xx + 0.16, 5.90, 3.92, 0.58, d, sz=6.0, c=SUB, leading=1.25)
        xx += 4.35
    rect(sl, 0.22, 6.68, 12.90, 0.52, L_PUR, B_PUR, 1.2, 0.07)
    txt(sl, 0.40, 6.68, 12.54, 0.52,
        [dict(runs=[('★ 若列级脱敏是硬需求：', {'b': True, 'c': PURD, 'sz': 7.2}),
                    ('B（Trino 读）的 UC 列掩码完全不可用（Iceberg REST Catalog 不支持），B2 原生生效、任何路径强制遮蔽——这一条基本决定必须选 B2',
                     {'c': '4C1D95', 'sz': 7.2})])], anchor='m')


# ============================================= SLIDE 6/7 · ER（draw.io 风格·两页） =
# 参考 draw.io ER 模板视觉：模块彩条 + 白卡彩框 + 表头色块 + PK/FK/UK 徽标 + 1/N 连线
ER_MODC = {
    'yellow': ('FFF2CC', 'D6B656', 'D6A22A'),
    'blue':   ('D9E8FC', 'A9C4E5', '5B82B1'),
    'green':  ('D5E8D4', '9BCB92', '82B366'),
    'purple': ('E1D5E7', 'B5A0C5', '9673A6'),
    'red':    ('F8CECC', 'D79B98', 'B85450'),
    'gray':   ('F5F5F5', 'A8A8A8', '666666'),
}
ER_MOD = {}
for n in ['user', 'user_role', 'role', 'role_permission', 'permission', 'page', 'role_page',
          'query_config', 'query_field', 'query_join', 'drill_config']:
    ER_MOD[n] = 'yellow'
for n in ['decommission_system', 'system_user', 'decommission_system†', 'sync_job',
          'sync_job_config', 'sync_table_config', 'sync_table_stat', 'archive_batch',
          'archive_file', 'archive_set', 'archive_set_item']:
    ER_MOD[n] = 'blue'
for n in ['source_database', 'source_table', 'unstructured_source', 'unstructured_item',
          'attachment_index']:
    ER_MOD[n] = 'green'
for n in ['schema', 'storage_config']:
    ER_MOD[n] = 'purple'
for n in ['audit_log', 'sync_activity', 'notification']:
    ER_MOD[n] = 'gray'
for n in ['retention_policy', 'retention_assignment', 'legal_hold_event', 'destroy_approval']:
    ER_MOD[n] = 'red'

# ---------- 上篇：平台 / 权限 / 数据源（18 表） ----------
ER_A = dict(
    title='元数据库 ER 图（上）· 平台 / 权限 / 数据源配置',
    note='18 表 · 仅关键字段 · 通用列已省略 · 源码：report/元数据库ER图.mmd',
    bands=[
        ('① 权限与访问控制（RBAC）', 0.30, 0.70, 11.20, 'yellow'),
        ('③ 源数据配置', 6.30, 2.10, 4.55, 'green'),
        ('④ 目标存储 · 归档 Schema', 11.10, 2.10, 1.95, 'purple'),
        ('② 退役系统管理（核心）', 2.90, 3.48, 2.90, 'blue'),
        ('⑩ 通知 · ⑪ 审计与监控', 0.30, 5.82, 12.75, 'gray'),
    ],
    entities=[
        ('user', 0.30, 1.10, 2.10, 1.20, [('bigint', 'id', 'PK'), ('varchar', 'username', 'UK'), ('varchar', 'password_hash', ''), ('varchar', 'real_name·email', ''), ('tinyint', 'status', ''), ('datetime', 'last_login_at', '')]),
        ('user_role', 2.65, 1.30, 1.85, 0.70, [('bigint', 'user_id', 'FK'), ('bigint', 'role_id', 'FK')]),
        ('role', 4.75, 1.10, 2.10, 0.90, [('bigint', 'id', 'PK'), ('varchar', 'role_code', 'UK'), ('varchar', 'role_name', ''), ('tinyint', 'is_builtin', '')]),
        ('role_permission', 7.10, 1.30, 1.95, 0.70, [('bigint', 'role_id', 'FK'), ('bigint', 'permission_id', 'FK')]),
        ('permission', 9.30, 1.10, 2.20, 0.90, [('bigint', 'id', 'PK'), ('varchar', 'permission_code', 'UK'), ('varchar', 'module', ''), ('varchar', 'resource_type', '')]),
        ('page', 0.30, 2.38, 2.10, 0.90, [('bigint', 'id', 'PK'), ('varchar', 'page_code', 'UK'), ('varchar', 'path', ''), ('tinyint', 'enabled', '')]),
        ('role_page', 2.65, 2.54, 1.85, 0.70, [('bigint', 'role_id', 'FK'), ('bigint', 'page_id', 'FK')]),
        ('decommission_system', 2.90, 3.80, 2.90, 1.35, [('bigint', 'id', 'PK'), ('varchar', 'system_code', 'UK'), ('varchar', 'status', ''), ('int', 'retention_years', ''), ('date', 'sync_completed_date', ''), ('date', 'destroy_after_date', '')]),
        ('system_user', 3.20, 5.28, 1.80, 0.70, [('bigint', 'system_id', 'FK'), ('bigint', 'user_id', 'FK')]),
        ('source_database', 6.30, 2.42, 2.15, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'db_type', ''), ('varchar', 'host:port', ''), ('varchar', 'credential_ref', '')]),
        ('source_table', 8.70, 2.42, 2.15, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'source_database_id', 'FK'), ('varchar', 'table_name', ''), ('varchar', 'primary_key', ''), ('tinyint', 'is_attachment', '')]),
        ('unstructured_source', 6.30, 3.76, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'source_type', ''), ('varchar', 'location_uri', '')]),
        ('unstructured_item', 8.70, 3.76, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'unstructured_source_id', 'FK'), ('varchar', 'original_path', ''), ('varchar', 'hash', '')]),
        ('schema', 11.10, 2.42, 1.95, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'catalog_name', ''), ('varchar', 'schema_name', '')]),
        ('storage_config', 11.10, 3.76, 1.95, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'storage_type', ''), ('varchar', 'container', '')]),
        ('audit_log', 0.30, 6.20, 2.60, 0.78, [('bigint', 'id', 'PK'), ('bigint', 'user_id', 'FK'), ('varchar', 'action·module', ''), ('varchar', 'ip·status', '')]),
        ('sync_activity', 3.06, 6.20, 2.86, 0.78, [('bigint', 'id', 'PK'), ('date', 'activity_date', ''), ('int', '成功/失败/运行', ''), ('', 'sync_job 按天预聚合·派生', '')]),
        ('notification', 6.08, 6.20, 2.60, 0.62, [('bigint', 'id', 'PK'), ('varchar', 'recipient', ''), ('varchar', 'channel', '')]),
    ],
    edges=[
        ([(2.40, 1.60), (2.65, 1.60)], None, False, 'yellow'),
        ([(4.75, 1.60), (4.50, 1.60)], None, False, 'yellow'),
        ([(6.85, 1.60), (7.10, 1.60)], None, False, 'yellow'),
        ([(9.30, 1.60), (9.05, 1.60)], None, False, 'yellow'),
        ([(2.40, 3.12), (2.65, 3.12)], None, False, 'yellow'),
        ([(5.20, 1.96), (5.20, 2.38), (3.57, 2.38), (3.57, 2.66)], None, False, 'yellow'),
        ([(2.40, 1.88), (2.55, 1.88), (2.55, 5.65), (3.20, 5.65)], None, False, 'blue'),
        ([(4.35, 5.40), (4.35, 5.56)], None, False, 'blue'),
        ([(5.80, 4.46), (6.05, 4.46), (6.05, 3.12), (6.30, 3.12)], None, False, 'green'),
        ([(5.80, 4.54), (6.30, 4.54)], None, False, 'green'),
        ([(8.45, 3.12), (8.70, 3.12)], None, False, 'green'),
        ([(8.45, 4.54), (8.70, 4.54)], None, False, 'green'),
        ([(5.80, 4.26), (6.15, 4.26), (6.15, 3.78), (10.98, 3.78), (10.98, 3.28), (11.10, 3.28)], None, False, 'purple'),
        ([(5.80, 5.02), (5.95, 5.02), (5.95, 5.18), (10.98, 5.18), (10.98, 4.54), (11.10, 4.54)], None, False, 'purple'),
        ([(0.30, 7.00), (0.18, 7.00), (0.18, 1.88), (0.30, 1.88)], None, True, 'gray'),
    ],
    chip=(4.85, 2.80, 1.18, 0.52, 'sync_job / query_config\n子域见下页 →'),
    chip2=(8.84, 6.26, 4.21, 0.58),
    extra=[('弱关联', 0.03, 3.55, 0.16, 0.62, True)],
)

# ---------- 下篇：同步归档 / 查询 / 保留销毁（17 表 + 跨页引用） ----------
ER_B = dict(
    title='元数据库 ER 图（下）· 同步归档 / 动态查询 / 保留销毁',
    note='17 表 + 1 跨页引用（虚线†见上页）\n记号：[PK]/[FK]/[UK] · 连线 1—N 为一对多',
    bands=[
        ('⑤ 同步与归档（含 ⑥ 附件索引）', 0.30, 0.58, 10.20, 'blue'),
        ('⑧ 生命周期与保留 · ⑨ 销毁审批', 10.55, 0.58, 2.10, 'red'),
        ('⑦ 动态查询', 0.30, 5.10, 10.20, 'yellow'),
    ],
    entities=[
        ('decommission_system†', 0.30, 0.96, 2.30, 0.70, [('bigint', 'id', 'PK'), ('varchar', 'system_code', 'UK')]),
        ('sync_job', 3.30, 0.96, 2.40, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('bigint', 'assignment_id', 'FK'), ('varchar', 'job_type', ''), ('varchar', 'status', '')]),
        ('sync_job_config', 6.10, 0.96, 2.05, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'cron_expr', ''), ('tinyint', 'enabled', '')]),
        ('retention_policy', 10.55, 0.96, 2.10, 1.00, [('bigint', 'id', 'PK'), ('varchar', 'policy_code', 'UK'), ('varchar', 'policy_name', ''), ('int', 'period_days', '')]),
        ('sync_table_config', 0.30, 2.35, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'source_table_id', 'FK'), ('tinyint', 'enabled', ''), ('int', 'retain_years', '')]),
        ('attachment_index', 2.75, 2.35, 2.30, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'table·business_key', ''), ('varchar', 'object_key', ''), ('varchar', 'sha256', '')]),
        ('sync_table_stat', 5.20, 2.35, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'sync_job_id', 'FK'), ('bigint', 'source_table_id', 'FK'), ('bigint', 'row_count', '')]),
        ('archive_batch', 7.65, 2.35, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'sync_job_id', 'FK'), ('int', 'batch_year', ''), ('bigint', 'rows_out', '')]),
        ('retention_assignment', 10.55, 2.35, 2.10, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'policy_id', 'FK'), ('varchar', 'object_type', ''), ('date', 'start·due_date', ''), ('varchar', 'status', '')]),
        ('archive_set_item', 2.75, 3.85, 2.30, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'archive_set_id', 'FK'), ('varchar', 'original_name', ''), ('varchar', 'blob_url', '')]),
        ('archive_set', 5.20, 3.85, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'archive_batch_id', 'FK'), ('varchar', 'set_name', ''), ('int', 'items·bytes', '')]),
        ('archive_file', 7.65, 3.85, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'archive_batch_id', 'FK'), ('varchar', 'blob_url', ''), ('varchar', 'checksum', '')]),
        ('legal_hold_event', 10.55, 3.85, 2.10, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'assignment_id', 'FK'), ('varchar', 'action', ''), ('varchar', 'reason', '')]),
        ('query_config', 0.30, 5.46, 2.40, 1.15, [('bigint', 'id', 'PK'), ('bigint', 'system_id', 'FK'), ('varchar', 'name', ''), ('varchar', 'base_table', ''), ('varchar', 'status', '')]),
        ('query_field', 3.30, 5.54, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'query_config_id', 'FK'), ('varchar', 'column_name', ''), ('tinyint', 'is_pii', '')]),
        ('query_join', 5.75, 5.54, 2.15, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'query_config_id', 'FK'), ('varchar', 'join_type', ''), ('varchar', 'left→right', '')]),
        ('drill_config', 8.20, 5.54, 2.30, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'query_config_id', 'FK'), ('bigint', 'schema_id', 'FK'), ('varchar', 'parent→child', '')]),
        ('destroy_approval', 10.55, 5.46, 2.10, 1.00, [('bigint', 'id', 'PK'), ('bigint', 'sync_job_id', 'FK'), ('bigint', 'approver_id', 'FK'), ('varchar', 'decision', '')]),
    ],
    edges=[
        ([(2.60, 1.35), (3.30, 1.35)], None, False, 'blue'),
        ([(1.45, 1.66), (1.45, 2.35)], None, False, 'blue'),
        ([(2.60, 1.60), (2.90, 1.60), (2.90, 2.02), (6.60, 2.02), (6.60, 1.74)], None, False, 'blue'),
        ([(2.60, 1.46), (2.67, 1.46), (2.67, 2.90), (2.75, 2.90)], None, False, 'green'),
        ([(0.30, 1.30), (0.18, 1.30), (0.18, 6.02), (0.30, 6.02)], None, False, 'yellow'),
        ([(5.45, 2.12), (5.45, 2.35)], None, False, 'blue'),
        ([(5.70, 1.74), (5.90, 1.74), (5.90, 2.12), (8.72, 2.12), (8.72, 2.35)], None, False, 'blue'),
        ([(10.55, 3.22), (10.30, 3.22), (10.30, 2.26), (5.65, 2.26), (5.65, 2.12)], None, False, 'red'),
        ([(5.70, 1.22), (5.90, 1.22), (5.90, 0.96), (13.00, 0.96), (13.00, 6.02), (12.65, 6.02)], None, False, 'red'),
        ([(8.72, 3.47), (8.72, 3.85)], None, False, 'blue'),
        ([(7.80, 3.47), (7.80, 3.70), (6.27, 3.70), (6.27, 3.85)], None, False, 'blue'),
        ([(5.20, 4.38), (5.05, 4.38)], None, False, 'blue'),
        ([(11.60, 1.74), (11.60, 2.35)], None, False, 'red'),
        ([(11.60, 3.47), (11.60, 3.85)], None, False, 'red'),
        ([(2.70, 6.15), (3.30, 6.15)], None, False, 'yellow'),
        ([(1.37, 6.58), (1.37, 6.80), (6.82, 6.80), (6.82, 6.50)], None, False, 'yellow'),
        ([(2.20, 6.58), (2.20, 6.98), (9.35, 6.98), (9.35, 6.50)], None, False, 'yellow'),
    ],
    chip=None,
    extra=[('到期触发销毁', 10.33, 2.32, 0.18, 0.56, True)],
)

def _er_card2(sl, name, x, y, w, h, rows):
    ghost = name.endswith('†')
    lf, bl, _ec = ER_MODC[ER_MOD[name]]
    rect(sl, x, y, w, h, 'FFFFFF', bl, 1.4, 0.04,
         dash=MSO_LINE_DASH_STYLE.DASH if ghost else None, shadow=not ghost)
    rect(sl, x + 0.014, y + 0.014, w - 0.028, 0.26, 'F5F5F5' if ghost else lf, None, 0, 0.03)
    txt(sl, x + 0.05, y + 0.045, w - 0.10, 0.22, name.replace('†', ''), sz=8.5, b=True,
        c='7A7A7A' if ghost else '222222', align='c', anchor='m')
    rect(sl, x + 0.014, y + 0.285, w - 0.028, h - 0.30, 'FAFAFA', None, 0, 0.02)
    paras = []
    tagc = {'PK': 'D00000', 'FK': '0049C6', 'UK': 'B26A00'}
    for tp, fn, tag in rows:
        runs = []
        if tag:
            runs.append((tag + '   ', {'sz': 6.4, 'b': True, 'c': tagc[tag]}))
        else:
            runs.append(('     ', {'sz': 6.4, 'c': 'FAFAFA'}))
        runs.append((fn, {'sz': 6.5, 'c': '222222'}))
        if tp:
            runs.append((' (' + tp + ')', {'sz': 6.0, 'c': '7A7A7A'}))
        paras.append(dict(runs=runs, ls=1.08))
    txt(sl, x + 0.12, y + 0.310, w - 0.20, h - 0.34, paras, leading=1.08)

def _er_edge2(sl, pts, color, lw=1.4):
    poly(sl, pts, color, lw, None)

def _er_1n(sl, pts):
    """父端 1 / 子端 N 标注（自动避线偏移）"""
    def unit(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        L = max(1e-9, (dx * dx + dy * dy) ** 0.5)
        return dx / L, dy / L
    ux, uy = unit(pts[0], pts[1])
    lx, ly = pts[0][0] + ux * 0.15 - uy * 0.13, pts[0][1] + uy * 0.15 + ux * 0.13
    txt(sl, lx - 0.06, ly - 0.065, 0.12, 0.13, '1', sz=7, b=True, c='222222', align='c', anchor='m')
    vx, vy = unit(pts[-1], pts[-2])
    rx, ry = pts[-1][0] + vx * 0.15 - vy * 0.13, pts[-1][1] + vy * 0.15 + vx * 0.13
    txt(sl, rx - 0.06, ry - 0.065, 0.12, 0.13, 'N', sz=7, b=True, c='222222', align='c', anchor='m')

def _er_slide2(prs, D, accent):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] += 1
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = _rgb(BG)
    title_bar(sl, D['title'], D['note'], accent)
    for label, bx, by, bw, bmod in D['bands']:
        lf, bl, _ec = ER_MODC[bmod]
        rect(sl, bx, by, bw, 0.38, lf, bl, 1.0, 0.06)
        txt(sl, bx + 0.12, by, bw - 0.24, 0.38, label, sz=9.5, b=True, c='333333', anchor='m')
    for ent in D['entities']:
        _er_card2(sl, *ent)
    for pts, label, dashed, mod in D['edges']:
        color = ER_MODC[mod][2]
        if dashed:
            for i in range(len(pts) - 1):
                seg(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], color, 1.2,
                    MSO_LINE_DASH_STYLE.DASH)
        else:
            _er_edge2(sl, pts, color)
        _er_1n(sl, pts)
        if label:
            txt(sl, label[1][0], label[1][1], 1.7, 0.12, label[0], sz=5.6, b=True, c=MUT)
    if D.get('chip2'):
        cx, cy, cw, chh = D['chip2']
        rect(sl, cx, cy, cw, chh, 'FFFFFF', 'A8A8A8', 1.0, 0.05, dash=MSO_LINE_DASH_STYLE.DASH)
        txt(sl, cx + 0.14, cy + 0.06, cw - 0.28, chh - 0.10,
            [dict(runs=[('跨页关联（详见卡内 FK）：', {'sz': 6.0, 'b': True, 'c': SLA})], ls=1.2),
             dict(runs=[('source_table → sync_table_config / sync_table_stat / archive_file', {'sz': 5.8, 'c': MUT})], ls=1.2),
             dict(runs=[('sync_job ⇢ sync_activity（按天预聚合，派生表无 FK）', {'sz': 5.8, 'c': MUT})], ls=1.2)])
    for ex in D.get('extra') or []:
        t, exx, exy, exw, exh, vert = ex
        txt(sl, exx, exy, exw, exh, t, sz=5.8, b=True, c=MUT if t == '弱关联' else ROS,
            align='c', anchor='m', vert=vert)
    if D.get('chip'):
        cx, cy, cw, chh, ct = D['chip']
        rect(sl, cx, cy, cw, chh, 'FFFFFF', 'B6AEE8', 1.0, 0.05, dash=MSO_LINE_DASH_STYLE.DASH)
        txt(sl, cx + 0.03, cy, cw - 0.06, chh, ct, sz=6.0, b=True, c=MUT, align='c', anchor='m', leading=1.25)
    # 底部图例
    _er_legend(sl, D['bands'])

def _er_band_bottom(D):
    """返回本 slide 内容区（非图例）的最低内容底部 y，用于定位图例。
    图例放在最低内容的底部下方，留 0.14" 间距。"""
    y = 0.70  # 内容区顶部基准
    for _, bx, by, bw, _ in D.get('bands', []):
        y = max(y, by + 0.38)  # 色带底部
    for ent in D.get('entities', []):
        ey = ent[3] + ent[4]  # y + h（底部边缘）
        y = max(y, ey)
    return y + 0.14  # 内容底部 + 0.14" 间距

def _er_legend(sl, bands):
    """页底图例：模块色块 + 连线标注，固定在幻灯片底部，避免与内容重叠。"""
    legend_y = 7.18  # 固定位置，紧贴幻灯片底部（7.5 - 0.32 = 7.18）
    rect(sl, 0.22, legend_y, 12.90, 0.32, 'F8FAFC', BORDER, 0.8, 0.04)
    txt(sl, 0.30, legend_y, 1.4, 0.32, '图例', sz=7.5, b=True, c=INK, anchor='m')
    # 模块色块
    xx = 0.80
    used = set()
    for label, bx, by, bw, bmod in bands:
        if bmod in used:
            continue
        used.add(bmod)
        lf, bl, ec = ER_MODC[bmod]
        oval(sl, xx, legend_y + 0.09, 0.14, 0.14, ec)
        txt(sl, xx + 0.20, legend_y, bw * 0.28, 0.32, label, sz=6.4, c=MUT, anchor='m')
        xx += bw * 0.28 + 0.18
    # 连线说明
    txt(sl, 10.05, legend_y, 3.0, 0.32,
        [dict(runs=[('1—N 连线：母端 1 / 子端 N　', {'c': SUB, 'sz': 6.2})], ls=1.15),
         dict(runs=[('（虚线卡=跨页引用）', {'c': FAINT, 'sz': 6.2})], ls=1.15)],
        align='r', anchor='m')

def slide_er(prs):
    _er_slide2(prs, ER_A, SLA)       # 第 6 页 · 上篇
    _er_slide2(prs, ER_B, PURD)      # 第 7 页 · 下篇

def emit_er_html(path):
    import os
    mmd_path = os.path.join(os.path.dirname(__file__), '元数据库ER图.mmd')
    mmd = open(mmd_path, encoding='utf-8').read()
    html = ('<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">'
            '<title>元数据库 ER 图 · 35 张表 / 11 模块</title>'
            '<style>body{margin:0;padding:28px;background:#f8fafc;font-family:\'Segoe UI\',\'PingFang SC\',\'Microsoft YaHei\',sans-serif;display:flex;flex-direction:column;align-items:center}'
            '.main{background:#fff;border:1px solid #e2e8f0;border-radius:14px;padding:24px;max-width:96vw;box-shadow:0 8px 30px rgba(15,23,42,.08)}'
            'h1{font-size:18px;color:#0f172a;margin:0 0 4px}p{font-size:12px;color:#64748b;margin:0 0 14px}'
            '.mermaid{display:flex;justify-content:center}</style>'
            '<script src="https://cdn.jsdelivr.net/npm/mermaid@11/dist/mermaid.min.js"></script>'
            '<script>mermaid.initialize({startOnLoad:true,theme:\"default\"});</script>'
            '</head><body><div class="main"><h1>RIMS 元数据库 ER 图 · 35 张表 / 11 模块</h1>'
            '<p>源码：report/元数据库ER图.mmd（Mermaid erDiagram）· 每表仅列关键字段，通用列已省略</p>'
            '<pre class="mermaid">\n' + mmd.replace('<', '&lt;').replace('&', '&amp;') + '</pre></div></body></html>')
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html)

# ------------------------------------------------------------------- build --
def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = '方案 B2 架构图集'
    prs.core_properties.author = 'RIMS 退役归档平台'
    slide_cover(prs)
    slide_azure(prs)
    slide_app(prs)
    slide_lake(prs)
    slide_enc(prs)
    slide_er(prs)
    import os as _os
    _report_dir = _os.path.dirname(__file__)
    prs.save(_os.path.join(_report_dir, '方案B_B2_架构图集.pptx'))
    emit_er_html(_os.path.join(_report_dir, '元数据库ER图.html'))
    with open('/tmp/arch_ops.json', 'w', encoding='utf-8') as f:
        json.dump(OPS, f, ensure_ascii=False)
    print('saved 方案B_B2_架构图集.pptx | shapes ops:', len(OPS))

if __name__ == '__main__':
    build('.')
