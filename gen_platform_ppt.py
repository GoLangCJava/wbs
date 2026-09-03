# -*- coding: utf-8 -*-
"""
生成《One-Abbott 共享应用平台 · 技术架构图集》PPTX
绘图助手与视觉体系复用 gen_arch_ppt.py（PPT 原生形状，不嵌图片）。
内容基于 archi.md（初步架构）与 tech-architecture.md（技术架构设计）。
同时输出 /tmp/arch_ops.json 供 PIL 预览渲染检查版式。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

from gen_arch_ppt import (
    INK, SUB, MUT, FAINT, BORDER, BG,
    AZ, AZD, GRN, GRND, PUR, PURD, CYN, CYND, ORG, AMB, AMBD, ROS, SLA,
    L_BLUE, B_BLUE, L_GRN, B_GRN, L_PUR, B_PUR, L_CYN, B_CYN,
    L_AMB, B_AMB, L_ROS, B_ROS, L_VNET, L_GRNB, B_GRNB, L_PURB, B_PURB, L_CYNB, B_CYNB,
    rect, oval, seg, poly, txt, chip, person, dbicon, cloud, lockicon, title_bar,
    OPS, _cur,
)

DASH = MSO_LINE_DASH_STYLE.DASH
OUT_NAME = 'One-Abbott_共享应用平台_技术架构图集.pptx'


# ============================================================== COVER =======
def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 0
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string('FFFFFF')
    rect(sl, 0, 0, 13.333, 0.12, AZ)
    rect(sl, 0, 0.12, 13.333, 0.035, GRN)
    rect(sl, 0, 0.155, 13.333, 0.035, PUR)
    rect(sl, 7.75, 1.30, 5.58, 5.9, 'F8FAFC', rad=0.25)
    rect(sl, 8.30, 1.85, 4.75, 4.75, 'FFFFFF', line_c=BORDER, lw=1.2, rad=0.22, shadow=True)

    # ---- 右侧分层速览 motif ----
    layers = [
        (2.30, L_GRN, B_GRN, GRND, '渠道层 · 微信 / 小程序 / 企微 / App / Portal'),
        (3.02, L_CYN, B_CYN, CYND, 'BFF / Portal Shell · Module Federation 微前端'),
        (3.74, L_AMB, B_AMB, AMBD, '能力层 · Content / Event / Customer / Sales / Service / AI / Plugin Hub'),
        (4.46, L_BLUE, B_BLUE, AZD, 'HCP Pool · OneID · Consent · Customer360（Fabric 湖仓）'),
        (5.18, L_PUR, B_PUR, PURD, 'AI 平台 · model-gateway · RAG · Azure OpenAI'),
    ]
    for i, (y0, lf, bl, tc, s) in enumerate(layers):
        rect(sl, 8.62, y0, 4.10, 0.52, lf, bl, 1.2, 0.08)
        txt(sl, 8.68, y0, 3.98, 0.52, s, sz=7.8, c=tc, b=True, align='c', anchor='m')
        if i < 4:
            seg(sl, 10.67, y0 + 0.52, 10.67, y0 + 0.70, MUT, 1.4)
    rect(sl, 8.62, 5.90, 4.10, 0.38, 'F6F8FA', 'CBD5E1', 1.1, 0.06)
    txt(sl, 8.68, 5.90, 3.98, 0.38, 'Azure Landing Zone · AKS · GitOps · 可观测', sz=7, b=True,
        c=SLA, align='c', anchor='m')
    txt(sl, 8.62, 6.34, 4.10, 0.24, 'Build Once · Reuse Everywhere', sz=7.5, c=FAINT, align='c', b=True)

    # ---- 左侧标题与目录 ----
    txt(sl, 0.9, 1.55, 6.6, 0.3, 'One-Abbott 共享应用平台', sz=11, c=MUT, spc=300, b=True)
    txt(sl, 0.88, 1.92, 6.9, 0.85, '技术架构图集', sz=33, b=True, c=INK)
    rect(sl, 0.92, 2.86, 0.62, 0.045, AZ)
    txt(sl, 0.92, 3.02, 6.4, 0.3, '总体架构 · 微前端与插件 · HCP Pool 数据架构 · AI 平台与合规 · 全部 PPT 原生形状绘制',
        sz=10.5, c=MUT)
    items = [
        ('01', AZ, L_BLUE, B_BLUE, '总体技术架构',
         '渠道 → 接入边缘 → BFF → 能力 Hub → 事件骨干 → 数据层 → Landing Zone ｜ AI 平台'),
        ('02', CYN, L_CYN, B_CYN, '微前端与插件架构',
         'Portal Shell · Module Federation · manifest 灰度发布 · 插件运行时与安全'),
        ('03', GRN, L_GRN, B_GRN, 'HCP Pool 数据架构',
         'OneID 匹配 · Consent 同意账本 · Customer360 · Fabric Medallion 湖仓'),
        ('04', PUR, L_PUR, B_PUR, 'AI 平台与医药合规',
         'model-gateway 场景路由 · RAG 管道 · AE / PQ 护栏 · 两步审核'),
    ]
    y = 3.42
    for no, col, lf, bl, t, d in items:
        rect(sl, 0.92, y, 6.35, 0.66, 'FFFFFF', BORDER, 1.2, 0.09, shadow=True)
        rect(sl, 1.06, y + 0.09, 0.48, 0.48, lf, bl, 1.2, 0.08)
        txt(sl, 1.06, y + 0.09, 0.48, 0.48, no, sz=12, b=True, c=col, align='c', anchor='m', mono=True)
        txt(sl, 1.68, y + 0.08, 5.5, 0.24, t, sz=10, b=True, c=INK)
        txt(sl, 1.68, y + 0.35, 5.5, 0.22, d, sz=7.0, c=MUT)
        y += 0.74
    rect(sl, 0.92, 6.98, 6.35, 0.035, BORDER)
    txt(sl, 0.92, 7.08, 6.4, 0.25, '整理自 archi.md 与 tech-architecture.md · 2026-09', sz=8, c=FAINT)


# ====================================================== SLIDE 1 · 总体架构 ==
def slide_overall(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 1
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string(BG)
    title_bar(sl, 'One-Abbott 共享应用平台 · 总体技术架构',
              'Build Once · Reuse Everywhere\n渠道 / 微前端 → 能力 Hub → HCP Pool → AI → Azure', AZ)

    # ---------- 01 渠道层 ----------
    rect(sl, 0.22, 0.62, 10.14, 0.76, L_GRN, B_GRN, 1.2, 0.08)
    txt(sl, 0.40, 0.69, 3.0, 0.17, '01 渠道层 · Experience', sz=8.5, b=True, c=GRND)
    txt(sl, 6.4, 0.70, 3.8, 0.14, 'HCP · 销售 · 服务工程师 · 内部运营', sz=6.2, c=MUT, align='r')
    chs = [
        ('微信公众号 · H5', '内容 / 会议 / 积分'),
        ('微信小程序 · Taro', '内容 / 签到 / AI 问答'),
        ('企业微信', '侧边栏 · 拜访 / C360'),
        ('App · React Native', '离线拜访 / 工单'),
        ('Portal · 微前端 Shell', '运营管理后台'),
    ]
    xx = 0.40
    for t, s in chs:
        chip(sl, xx, 0.90, 1.86, 0.42, t, s, 'FFFFFF', 'BBF7D0', 1.0, 0.05,
             tsz=7.4, sc=MUT, ssz=6.0, shadow=True)
        xx += 1.955
    for ax in (2.55, 7.20):
        seg(sl, ax, 1.38, ax, 1.46, GRN, 1.8)

    # ---------- 02 接入边缘 ----------
    rect(sl, 0.22, 1.46, 10.14, 0.56, 'F8FAFC', 'CBD5E1', 1.2, 0.07)
    txt(sl, 0.40, 1.53, 2.6, 0.15, '02 接入边缘', sz=8.5, b=True, c=SUB)
    chip(sl, 0.40, 1.72, 2.10, 0.24, 'CDN + WAF', None, 'FFFFFF', 'CBD5E1', 1.0, 0.04, tc=INK, tsz=6.8)
    seg(sl, 2.52, 1.84, 2.72, 1.84, MUT, 1.4)
    chip(sl, 2.74, 1.72, 3.30, 0.24, 'Azure API Management · 外部实例', None,
         'FFFFFF', AZ, 1.2, 0.04, tc=AZD, tsz=6.8)
    chip(sl, 6.20, 1.66, 3.90, 0.36, '身份：Entra ID（内部 SSO·MFA）',
         'identity-service 微信 OAuth / 验证码（外部 → OneID）',
         'FFFFFF', B_BLUE, 1.0, 0.05, tc=INK, tsz=6.8, sc=MUT, ssz=5.8)
    for ax in (2.55, 7.20):
        seg(sl, ax, 2.02, ax, 2.10, CYN, 1.8)

    # ---------- 03 BFF ----------
    rect(sl, 0.22, 2.10, 10.14, 0.64, L_CYN, B_CYNB, 1.2, 0.08)
    txt(sl, 0.40, 2.17, 3.4, 0.15, '03 BFF 聚合层 · AKS', sz=8.5, b=True, c=CYND)
    txt(sl, 5.4, 2.18, 4.8, 0.13, '只做聚合 / 裁剪 / 鉴权上下文 · 不写业务规则', sz=6.2, c=MUT, align='r')
    bfs = [('hcp-bff', '小程序 / H5 聚合'), ('wework-bff', '企微侧边栏聚合'),
           ('app-bff', '离线同步 · 冲突合并'), ('portal-bff', '运营门户聚合')]
    xx = 0.40
    for t, s in bfs:
        chip(sl, xx, 2.32, 2.36, 0.38, t, s, 'FFFFFF', CYN, 1.2, 0.05, tsz=7.6, sc=MUT, ssz=6.0)
        xx += 2.44
    for ax in (2.55, 7.20):
        seg(sl, ax, 2.74, ax, 2.82, CYN, 1.8)

    # ---------- 04 能力层 ----------
    rect(sl, 0.22, 2.82, 10.14, 2.04, 'FFFFFF', BORDER, 1.2, 0.08)
    txt(sl, 0.40, 2.89, 7.2, 0.15, '04 能力层 · Capability Hub（AKS 命名空间 · 一 Hub 一限界上下文 · Database-per-Service）',
        sz=8.5, b=True, c='334155')
    rect(sl, 0.36, 3.10, 0.34, 1.44, L_BLUE, B_BLUE, 1.2, 0.05)
    txt(sl, 0.36, 3.10, 0.34, 1.44, '内部 APIM · VNet 内', sz=6.0, b=True, c=AZD,
        align='c', anchor='m', vert=True)
    seg(sl, 0.70, 3.41, 0.82, 3.41, AZ, 1.4)
    seg(sl, 0.70, 4.23, 0.82, 4.23, AZ, 1.4)
    hubs_r1 = [
        ('Content Hub · 内容与学术', 'content · dam · search · publish', 'CMS · DAM · 多渠道发布', L_BLUE, B_BLUE, AZD),
        ('Event Hub · 会议与积分', 'event · registration · live · points', '签到 · 直播 · 双录积分账本', L_GRN, B_GRN, GRND),
        ('Customer Hub · 客户主数据', 'hcp/hco · oneid · consent · c360 · tag', '匹配合并 · 同意账本 · 标签', L_CYN, B_CYN, CYND),
        ('Sales Hub · 销售执行', 'visit · checkin · sync · performance', 'GPS 打卡 · 离线同步 · KPI', L_AMB, B_AMB, AMBD),
    ]
    hubs_r2 = [
        ('Service Hub · 售后服务', 'ticket · parts · iot-ingest', 'SLA 计时 · 备件 · 遥测告警', 'F6F8FA', 'CBD5E1', SLA),
        ('Plugin Runtime · 插件', '术后随访 · 医学咨询 · 竞品情报 · 差旅', '远程 MF ＋ 插件 API · HMAC · 租户启用', 'FDF2F8', 'F9A8D4', 'BE185D'),
        ('横切服务 · Cross-cutting', 'identity · notification · workflow · rules', 'file · audit · 契约 @abbott/contracts', 'FFF7ED', 'FED7AA', AMBD),
        ('AI Hub · 智能服务', 'copilot · recommend · review', '模型网关 ＋ RAG → 见右侧 AI 平台', L_PUR, B_PUR, PURD),
    ]
    for row, hubs, y0 in ((1, hubs_r1, 3.04), (2, hubs_r2, 3.86)):
        xx = 0.82
        for t, s1, s2, lf, bl, tc in hubs:
            rect(sl, xx, y0, 2.26, 0.74, lf, bl, 1.2, 0.06, shadow=True)
            txt(sl, xx + 0.06, y0 + 0.05, 2.14, 0.15, t, sz=7.2, b=True, c=INK, align='c')
            txt(sl, xx + 0.06, y0 + 0.24, 2.14, 0.13, s1, sz=5.8, c=tc, align='c', b=True)
            txt(sl, xx + 0.06, y0 + 0.40, 2.14, 0.28, s2, sz=5.6, c=MUT, align='c', leading=1.1)
            xx += 2.34
    for ax in (2.55, 7.20):
        seg(sl, ax, 4.86, ax, 4.94, AMB, 1.8)

    # ---------- 05 事件骨干 ----------
    rect(sl, 0.22, 4.94, 10.14, 0.46, 'FFF7ED', 'FED7AA', 1.2, 0.07)
    txt(sl, 0.40, 5.00, 1.4, 0.14, '05 事件骨干', sz=7.8, b=True, c=AMBD)
    txt(sl, 0.40, 5.16, 1.4, 0.12, 'EVENT DRIVEN', sz=5.2, b=True, c=AMB, spc=110, mono=True)
    chip(sl, 1.85, 5.02, 3.95, 0.30,
         [dict(runs=[('Service Bus · 领域事件', {'c': AMBD, 'sz': 6.6, 'b': True}),
                     ('　CloudEvents · Outbox · 幂等 · DLQ · 对账', {'c': MUT, 'sz': 6.2})])],
         None, 'FFFFFF', ORG, 1.2, 0.05)
    chip(sl, 5.94, 5.02, 4.20, 0.30,
         [dict(runs=[('Event Hub · 行为流', {'c': AMBD, 'sz': 6.6, 'b': True}),
                     ('　埋点采集 → 实时标签(KEDA) → Fabric', {'c': MUT, 'sz': 6.2})])],
         None, 'FFFFFF', ORG, 1.2, 0.05)
    for ax in (2.55, 9.20):
        seg(sl, ax, 5.40, ax, 5.48, AZ, 1.8)

    # ---------- 06 数据层 ----------
    rect(sl, 0.22, 5.48, 10.14, 0.84, L_BLUE, B_BLUE, 1.2, 0.08)
    txt(sl, 0.40, 5.55, 6.0, 0.15, '06 数据层 · HCP Pool 与湖仓（Medallion）', sz=8.5, b=True, c='1D4ED8')
    stores = [
        ('Azure SQL', '每服务一库 · RLS 多租户'), ('Cosmos DB', '行为 · 画像 · 会话'),
        ('Redis', '缓存 · 去重 · 限流'), ('Blob / ADLS', '资产 · 归档 · WORM'),
        ('Fabric OneLake', 'B→S→G · C360 宽表'), ('Databricks', 'OneID 匹配 · ML'),
        ('Purview', '分级 · 血缘 · 术语'),
    ]
    xx = 0.40
    for t, s in stores:
        chip(sl, xx, 5.74, 1.32, 0.48, t, s, 'FFFFFF', '93C5FD', 1.0, 0.05, tc=INK, tsz=6.8, sc=MUT, ssz=5.6)
        xx += 1.398

    # ---------- 07 底座 ----------
    rect(sl, 0.22, 6.40, 10.14, 0.56, 'F6F8FA', 'CBD5E1', 1.2, 0.07)
    txt(sl, 0.40, 6.46, 5.0, 0.14, '07 Azure Landing Zone · 平台工程', sz=8, b=True, c=SLA)
    plats = ['AKS 私有集群 · KEDA 弹性', 'Terraform + ArgoCD · GitOps 金丝雀',
             'Monitor + Grafana · OTel / SLO', 'Key Vault · App Config · CMK / Flag']
    xx = 0.40
    for s in plats:
        chip(sl, xx, 6.63, 2.40, 0.26, s, None, 'FFFFFF', 'CBD5E1', 1.0, 0.04, tc=SUB, tsz=6.2)
        xx += 2.48

    # ---------- 右侧 AI 平台 ----------
    rect(sl, 10.50, 0.62, 2.61, 5.68, 'FAF9FF', PUR, 1.4, 0.09)
    txt(sl, 10.64, 0.72, 2.35, 0.18, 'AI 平台 · AI Layer', sz=9.5, b=True, c=PUR)
    txt(sl, 10.64, 0.92, 2.35, 0.14, 'Azure OpenAI ＋ 国产备案模型 · 场景路由', sz=6.0, c=MUT)
    ai_cards = [
        (1.12, 0.66, L_PUR, PUR, 1.5, 'model-gateway 统一网关',
         '场景路由 · 配额成本 · 审计 · 降级', '一切模型调用必经 ★'),
        (1.86, 0.54, 'FFFFFF', B_PUR, 1.1, 'rag-service 检索增强',
         'hybrid 检索 · 行级权限 · 引用必附', None),
        (2.48, 0.54, 'FFFFFF', B_PUR, 1.1, 'copilot-api 会话编排',
         'HCP 助手 · 销售 / 运营 Copilot', None),
        (3.10, 0.54, 'FFFFFF', B_PUR, 1.1, 'recommend · review',
         '个性化推荐 · AI 合规预审', None),
    ]
    for y0, h0, lf, bl, lw0, t, s1, s2 in ai_cards:
        rect(sl, 10.64, y0, 2.33, h0, lf, bl, lw0, 0.05, shadow=True)
        paras = [(t, {'b': True, 'sz': 7.2, 'c': INK}), (s1, {'sz': 5.8, 'c': MUT})]
        if s2:
            paras.append((s2, {'sz': 5.7, 'c': PUR, 'b': True}))
        txt(sl, 10.70, y0 + 0.04, 2.21, h0 - 0.08, paras, align='c', anchor='m', leading=1.15)
    chip(sl, 10.64, 3.72, 1.13, 0.28, 'AI Search', None, L_CYN, B_CYNB, 1.0, 0.04, tc=CYND, tsz=6.0)
    chip(sl, 11.84, 3.72, 1.13, 0.28, 'Content Safety', None, L_CYN, B_CYNB, 1.0, 0.04, tc=CYND, tsz=5.7)
    chip(sl, 10.64, 4.06, 2.33, 0.28, 'Prompt Flow 评测 · 入 CI 门禁', None, 'FFFFFF', BORDER, 1.0, 0.04,
         tc=SUB, tsz=6.0)
    rect(sl, 10.64, 4.44, 2.33, 0.62, L_ROS, ROS, 1.3, 0.05, dash=DASH)
    txt(sl, 10.70, 4.48, 2.21, 0.54,
        [('AE / PQ 出站护栏 ★', {'b': True, 'sz': 7.0, 'c': ROS}),
         ('自动捕获 → PV 工单 · 会话全留痕', {'sz': 5.7, 'c': SUB})], align='c', anchor='m', leading=1.15)
    rect(sl, 10.64, 5.14, 2.33, 0.58, 'FFFFFF', BORDER, 1.0, 0.05)
    txt(sl, 10.70, 5.18, 2.21, 0.50,
        [('模型路线 C（推荐）', {'b': True, 'sz': 6.4, 'c': INK}),
         ('非敏 → AOAI · 敏感 → 境内驻留', {'sz': 5.7, 'c': MUT})], align='c', anchor='m', leading=1.15)
    txt(sl, 10.64, 5.80, 2.33, 0.30, 'Phase 3：AI Search → 推荐 →\nCopilot → AI 合规审核', sz=5.8, c=FAINT,
        align='c', leading=1.2)
    seg(sl, 10.12, 4.23, 10.48, 4.23, PUR, 1.6, DASH)

    # ---------- 底部说明 ----------
    txt(sl, 0.22, 7.04, 12.9, 0.16,
        '分层对应 archi.md：Experience → Application → Capability → Data → AI → Cloud ｜ 多租户：tid 全链路 ＋ SQL 行级安全 ｜ 发布：GitOps 金丝雀 ＋ 微前端 manifest 灰度',
        sz=6.4, c=MUT, align='c')


# ================================================= SLIDE 2 · 微前端与插件 ==
def slide_mfe(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 2
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string(BG)
    title_bar(sl, '微前端与插件架构 · Micro Frontend & Plugin Runtime',
              'React 18 + TS5 + Rspack/Vite + Module Federation\n小程序 Taro · App React Native · 多端复用', CYN)

    # ---- 左上：Portal Shell ----
    rect(sl, 0.22, 0.62, 6.30, 1.30, 'FFFFFF', BORDER, 1.2, 0.07, shadow=True)
    txt(sl, 0.36, 0.70, 5.9, 0.16, 'Portal Shell · Next.js Host（平台团队统一维护）', sz=8.2, b=True, c=INK)
    shell = ['路由 / 布局\nRBAC 菜单下发', 'SSO 会话\n权限 SDK', 'Design Tokens\n主题 / 暗色', '事件总线\nmitt', 'MF 注册表\nmanifest 灰度']
    xx = 0.36
    for s in shell:
        chip(sl, xx, 0.92, 1.14, 0.48, s, None, L_CYN, B_CYNB, 1.0, 0.05, tc=CYND, tsz=6.0, leading=1.1)
        xx += 1.21
    txt(sl, 0.36, 1.50, 6.0, 0.14, '全局错误兜底：模块加载失败 → 降级占位不白屏 · 菜单 / 权限按 RBAC 动态下发',
        sz=6.0, c=MUT)
    # Shell -> MFs
    seg(sl, 3.37, 1.92, 3.37, 2.02, CYN, 1.8)

    # ---- 左中：远程模块 ----
    rect(sl, 0.22, 2.02, 6.30, 1.56, L_CYN, B_CYNB, 1.2, 0.07)
    txt(sl, 0.36, 2.09, 5.0, 0.15, '远程模块 Remote MF · 独立仓库 / 独立流水线（8 个）', sz=8, b=True, c=CYND)
    mfs = [('产品', 'product-mf'), ('内容学术', 'content-mf'), ('会议', 'event-mf'), ('销售', 'sales-mf'),
           ('售后', 'service-mf'), ('客户', 'customer-mf'), ('AI 助手', 'ai-mf'), ('插件', 'plugin-mf')]
    for i, (t, s) in enumerate(mfs):
        mx = 0.36 + (i % 4) * 1.525; my = 2.30 + (i // 4) * 0.62
        chip(sl, mx, my, 1.44, 0.54, t, s, 'FFFFFF', CYN, 1.2, 0.05, tc=INK, tsz=7.4, sc=MUT, ssz=5.9)
    txt(sl, 0.36, 3.42, 6.0, 0.12, '跨模块通信仅走事件总线 mf:<module>:<event> · 禁止跨 MF 直接 import',
        sz=6.0, b=True, c=CYND)

    # ---- 左下：共享与发布 ----
    rect(sl, 0.22, 3.68, 6.30, 0.40, L_AMB, B_AMB, 1.1, 0.06)
    txt(sl, 0.36, 3.68, 6.02, 0.40,
        [dict(runs=[('shared 单例白名单：', {'b': True, 'c': AMBD, 'sz': 6.4}),
                    ('react · design-system · auth-sdk · api-client · i18n · telemetry（白名单外独立打包，杜绝版本地狱）',
                     {'c': SUB, 'sz': 6.4})])], anchor='m')
    rect(sl, 0.22, 4.18, 6.30, 0.94, 'FFFFFF', BORDER, 1.2, 0.07)
    txt(sl, 0.36, 4.25, 5.9, 0.14, '发布与灰度（回滚 = manifest 指回旧版，分钟级）', sz=7.6, b=True, c=INK)
    steps = ['① 构建\nCI 产物', '② CDN\n不可变版本', '③ manifest\nPR 审批', '④ 灰度\n1→10→100%', '⑤ 异常\n自动回滚']
    xx = 0.36
    for i, s in enumerate(steps):
        chip(sl, xx, 4.46, 1.08, 0.50, s, None, L_GRN, B_GRN, 1.1, 0.05, tc=GRND, tsz=6.0, leading=1.12)
        if i < 4:
            seg(sl, xx + 1.08, 4.71, xx + 1.22, 4.71, GRN, 1.4)
        xx += 1.22
    rect(sl, 0.22, 5.22, 6.30, 0.40, L_GRN, B_GRN, 1.1, 0.06)
    txt(sl, 0.36, 5.22, 6.02, 0.40,
        [dict(runs=[('契约与门禁：', {'b': True, 'c': GRND, 'sz': 6.4}),
                    ('@abbott/contracts（API / 事件 / 权限码同源）· Playwright 冒烟 ＋ Shell E2E · manifest 兼容性 CI',
                     {'c': SUB, 'sz': 6.4})])], anchor='m')

    # ---- 左底：渠道适配 ----
    rect(sl, 0.22, 5.72, 6.30, 0.62, 'F8FAFC', 'CBD5E1', 1.1, 0.06)
    txt(sl, 0.36, 5.78, 5.0, 0.13, '多端适配（同一套能力，按渠道裁剪）', sz=7.2, b=True, c=SUB)
    for i, s in enumerate(['小程序 · Taro 分包', '企微 / 公众号 H5', 'App · React Native', '共享逻辑层 Monorepo']):
        chip(sl, 0.36 + i * 1.51, 5.96, 1.44, 0.30, s, None, 'FFFFFF', 'CBD5E1', 1.0, 0.04, tc=SUB, tsz=6.0)

    # ---- 右列：插件运行时 ----
    RX = 6.66; RW = 6.45
    rect(sl, RX, 0.62, RW, 0.74, L_PUR, B_PUR, 1.3, 0.07)
    txt(sl, RX + 0.14, 0.70, RW - 0.28, 0.60,
        [('Plugin Hub · 插件运行时', {'b': True, 'sz': 8.4, 'c': PURD}),
         ('插件 = 远程 MF ＋ 后端插件 API ＋ manifest 声明（菜单挂点 · 权限 scope · 订阅事件 · 回调 URL）',
          {'sz': 6.2, 'c': SUB})], leading=1.25)
    # 生命周期
    txt(sl, RX, 1.48, 3.0, 0.14, '插件生命周期', sz=8, b=True, c=INK)
    life = [('① 注册', '安全评审 · 契约校验'), ('② 启用', '租户级 feature flag'),
            ('③ 运行', '失败自动降级隔离'), ('④ 下线', '数据导出 ＋ 归档')]
    xx = RX
    for i, (t, s) in enumerate(life):
        chip(sl, xx, 1.68, 1.50, 0.56, t, s, 'FFFFFF', B_PUR, 1.2, 0.05, tc=INK, tsz=7.2, sc=MUT, ssz=5.9)
        if i < 3:
            seg(sl, xx + 1.50, 1.96, xx + 1.61, 1.96, PUR, 1.4)
        xx += 1.61
    # 安全
    rect(sl, RX, 2.36, RW, 0.56, L_BLUE, B_BLUE, 1.1, 0.06)
    txt(sl, RX + 0.14, 2.36, RW - 0.28, 0.56,
        [dict(runs=[('插件 API 安全：', {'b': True, 'c': AZD, 'sz': 6.6}),
                    ('统一经 APIM plugin product 暴露 · HMAC 签名 ＋ IP 白名单 · 限流配额 · 不能直连内网',
                     {'c': SUB, 'sz': 6.6})])], anchor='m')
    # 插件目录
    rect(sl, RX, 3.02, RW, 1.34, 'FFFFFF', BORDER, 1.2, 0.07)
    txt(sl, RX + 0.14, 3.09, 4.0, 0.14, '插件目录（首期场景）', sz=8, b=True, c=INK)
    plugs = ['术后随访', '医学咨询', '竞品情报', '成本计算', '差旅管理', '＋ 第三方扩展']
    for i, s in enumerate(plugs):
        px = RX + 0.14 + (i % 3) * 2.10; py = 3.30 + (i // 3) * 0.50
        dash = True if i == 5 else None
        chip(sl, px, py, 1.98, 0.42, s, None,
             'FDF2F8' if i == 5 else 'FFFFFF', 'F9A8D4' if i == 5 else B_PUR,
             1.1, 0.05, tc='BE185D' if i == 5 else INK, tsz=7.0, dash=dash)
    # 价值
    rect(sl, RX, 4.46, RW, 0.62, L_GRN, B_GRN, 1.1, 0.06)
    txt(sl, RX + 0.14, 4.46, RW - 0.28, 0.62,
        [dict(runs=[('Build Once, Reuse Everywhere：', {'b': True, 'c': GRND, 'sz': 6.8}),
                    ('新场景 = 挂载 MF ＋ 注册 API，不动平台内核；租户按需启用、按标签计量计费',
                     {'c': SUB, 'sz': 6.8})])], anchor='m', leading=1.2)
    # 治理
    rect(sl, RX, 5.18, RW, 1.16, 'FFF7ED', 'FDBA74', 1.1, 0.06)
    txt(sl, RX + 0.14, 5.25, 3.0, 0.14, '微前端 / 插件治理（防腐化）', sz=7.8, b=True, c='B45309')
    govs = ['shared 白名单 ＋ 版本范围（react ^18 singleton）· CI 门禁校验兼容',
            '样式隔离：Design Tokens ＋ 模块级 class 前缀 · 全局样式仅 Shell 注入一次',
            '契约破坏即阻断合并 · 季度架构体检（依赖漂移 / 跨模块引用扫描）']
    for i, s in enumerate(govs):
        txt(sl, RX + 0.14, 5.44 + i * 0.26, RW - 0.28, 0.24,
            [dict(runs=[['· ', {'c': AMB, 'b': True, 'sz': 6.2}], [s, {'c': SUB, 'sz': 6.2}]], ls=1.15)])
    # 底部
    rect(sl, 0.22, 6.48, 12.89, 0.42, 'FFFFFF', 'CBD5E1', 1.0, 0.06)
    txt(sl, 0.40, 6.48, 12.5, 0.42,
        [dict(runs=[('技术栈：', {'b': True, 'c': '334155', 'sz': 6.6}),
                    ('Module Federation 运行时远程加载（single-spa 仅作异构框架逃生舱）｜ Portal Next.js 14 · HCP 端纯 SPA ＋ CDN 分发 ｜ 微信审核发布流水线（体验包 → 审核 → 发布）',
                     {'c': SUB, 'sz': 6.6})])], anchor='m')
    # 图例
    leg = [('Portal Shell（平台团队）', SLA), ('远程模块 / 插件', PUR), ('发布链路', GRN), ('安全 / 契约', AZ)]
    xx = 0.40
    for t, c in leg:
        oval(sl, xx, 7.06, 0.09, 0.09, c)
        txt(sl, xx + 0.14, 7.01, len(t) * 0.10 + 0.30, 0.16, t, sz=6.4, c=MUT)
        xx += 0.32 + len(t) * 0.10 + 0.50


# =================================================== SLIDE 3 · HCP POOL ====
def slide_hcp(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 3
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string(BG)
    title_bar(sl, 'HCP Pool · 数据架构（OneID · Consent · Customer360）',
              '数据流：Channel → Behavior Event → HCP Pool →\nCustomer360 → AI Platform（archi.md §9）', GRN)

    # ---- 左：Medallion 管道 ----
    stages = [
        ('01', CYN, L_CYN, B_CYNB, CYND, '渠道行为采集', '埋点 SDK → collector → Event Hub（分区按租户）'),
        ('02', SLA, 'F6F8FA', 'CBD5E1', SLA, 'Bronze 原始层', '存量四系统快照 · SQL CDC Mirroring · 行为流 Capture'),
        ('03', AZ, L_BLUE, B_BLUE, AZD, 'Silver 标准层', '清洗标准化 · OneID 关联 · ★Consent 过滤 · 脱敏'),
        ('04', PUR, L_PUR, B_PUR, PURD, 'Gold 指标层', '维度 / 事实 · Customer360 宽表 · 销售 KPI'),
        ('05', GRN, L_GRN, B_GRN, GRND, 'Serving 服务层', 'Power BI 看板 · c360-api（Redis 热点 ＋ Cosmos 全量）'),
    ]
    y0 = 0.62
    for i, (sid, col, lf, bl, tc, t, s) in enumerate(stages):
        rect(sl, 0.22, y0, 4.60, 0.60, lf, bl, 1.2, 0.07)
        oval(sl, 0.34, y0 + 0.13, 0.34, 0.34, 'FFFFFF', col, 1.6)
        txt(sl, 0.34, y0 + 0.13, 0.34, 0.34, sid, sz=6.5, b=True, c=col, align='c', anchor='m', mono=True)
        txt(sl, 0.78, y0 + 0.07, 3.9, 0.16, t, sz=8.4, b=True, c=INK)
        txt(sl, 0.78, y0 + 0.28, 3.94, 0.26, s, sz=6.0, c=MUT, leading=1.12)
        if i < 4:
            seg(sl, 2.52, y0 + 0.60, 2.52, y0 + 0.70, MUT, 1.6)
        y0 += 0.70
    rect(sl, 0.22, 4.14, 4.60, 0.40, 'F0FAF6', GRN, 1.1, 0.06)
    txt(sl, 0.36, 4.14, 4.34, 0.40,
        [dict(runs=[('存储与计算：', {'b': True, 'c': GRND, 'sz': 6.4}),
                    ('OneLake（Delta）· Databricks 匹配 / 训练 · Cosmos 画像 · Purview 血缘',
                     {'c': SUB, 'sz': 6.4})])], anchor='m')

    # ---- 中：OneID 匹配 ----
    MX = 5.00; MW = 3.92
    rect(sl, MX, 0.62, MW, 0.34, '334155', None, 0, 0.05)
    txt(sl, MX, 0.62, MW, 0.34, 'OneID 身份归一（平台级唯一标识 · 永不变更）', sz=7.6, b=True, c='FFFFFF',
        align='c', anchor='m')
    chip(sl, MX, 1.04, MW, 0.34, '来源：渠道注册 · 存量导入 · 第三方参考库（需法务评估来源合法性）',
         None, 'FFFFFF', BORDER, 1.0, 0.04, tc=SUB, tsz=6.2)
    seg(sl, MX + MW / 2, 1.38, MX + MW / 2, 1.46, MUT, 1.4)
    chip(sl, MX, 1.46, MW, 0.34, '标准化：姓名拼音 · 医院别名归一 · 科室字典', None,
         L_CYN, B_CYNB, 1.1, 0.04, tc=CYND, tsz=6.4)
    seg(sl, MX + MW / 2, 1.80, MX + MW / 2, 1.88, MUT, 1.4)
    match = [
        ('确定性匹配', '医师编号 / 证件 / 手机 Hash → 自动关联', L_GRN, B_GRN, GRND),
        ('概率匹配 p ≥ 0.95', '加权相似度（姓名＋医院＋科室＋职称）→ 自动合并', L_GRN, B_GRN, GRND),
        ('0.70 ≤ p < 0.95', 'Steward 人工审核队列（确认合并 / 否决新建）', L_AMB, B_AMB, AMBD),
        ('p < 0.70', '新建 OneID', 'F6F8FA', 'CBD5E1', SLA),
    ]
    yy = 1.88
    for t, s, lf, bl, tc in match:
        chip(sl, MX, yy, MW, 0.36, t, s, lf, bl, 1.2, 0.05, tc=INK, tsz=6.8, sc=MUT, ssz=5.9)
        seg(sl, MX + MW / 2, yy + 0.36, MX + MW / 2, yy + 0.44, MUT, 1.2)
        yy += 0.44
    rect(sl, MX, 3.66, MW, 0.50, L_ROS, ROS, 1.1, 0.05, dash=DASH)
    txt(sl, MX + 0.10, 3.66, MW - 0.20, 0.50,
        [('合并 / 拆分：全程审计 · 不可逆操作二次确认', {'b': True, 'sz': 6.4, 'c': ROS}),
         ('人工审核结果每月回流再训练匹配模型', {'sz': 5.9, 'c': SUB})], anchor='m', leading=1.15)
    chip(sl, MX, 4.24, MW, 0.30, 'OneID ↔ 微信 UnionID / 手机号 / 企微 / 存量系统 ID（一对多渠道身份）',
         None, 'FFFFFF', BORDER, 1.0, 0.04, tc=SUB, tsz=6.0)

    # ---- 右：Consent / 标签 / C360 / 治理 ----
    RX = 9.06; RW = 4.05
    rect(sl, RX, 0.62, RW, 1.66, L_GRN, B_GRN, 1.3, 0.07)
    txt(sl, RX + 0.12, 0.69, RW - 0.24, 0.15, 'Consent 同意账本（PIPL 合规核心）', sz=8.2, b=True, c=GRND)
    scopes = ['marketing', 'event_invite', 'profiling', 'ai_personalization', 'third_party_share']
    xx, srow = RX + 0.12, 0
    for s in scopes:
        w = 0.16 + len(s) * 0.052
        if xx + w > RX + RW - 0.12:
            xx = RX + 0.12; srow += 1
        chip(sl, xx, 0.88 + srow * 0.26, w, 0.22, s, None, 'FFFFFF', 'BBF7D0', 1.0, 0.04, tc=GRND, tsz=5.8)
        xx += w + 0.06
    txt(sl, RX + 0.12, 1.42, RW - 0.24, 0.44,
        [('采集：注册 / 活动表单 / 小程序授权（存证据与版本）', {'sz': 6.0, 'c': SUB}),
         ('触达前置校验 API：notification / 推荐调用前强制', {'sz': 6.0, 'c': SUB})], leading=1.25)
    rect(sl, RX + 0.12, 1.92, RW - 0.24, 0.28, 'FFFFFF', 'BBF7D0', 1.0, 0.04)
    txt(sl, RX + 0.12, 1.92, RW - 0.24, 0.28,
        [dict(runs=[('撤回 ≤ 24h 全渠道生效', {'b': True, 'c': GRND, 'sz': 6.4}),
                    ('（consent.revoked 广播 ＋ 缓存 TTL）', {'c': MUT, 'sz': 5.9})])], align='c', anchor='m')

    rect(sl, RX, 2.40, RW, 0.88, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, RX + 0.12, 2.47, RW - 0.24, 0.14, '标签引擎', sz=8, b=True, c=INK)
    txt(sl, RX + 0.12, 2.64, RW - 0.24, 0.58,
        [('规则标签（rules-engine 表达式）· 统计标签（日批 Fabric）', {'sz': 6.0, 'c': SUB}),
         ('实时标签（KEDA 消费行为流，如「近 1h 观看直播中」）', {'sz': 6.0, 'c': SUB}),
         ('字典治理：owner · 口径 · 生命周期（Purview 术语库）', {'sz': 6.0, 'c': MUT})], leading=1.22)

    rect(sl, RX, 3.38, RW, 0.92, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, RX + 0.12, 3.45, RW - 0.24, 0.14, 'Customer360 服务', sz=8, b=True, c=INK)
    txt(sl, RX + 0.12, 3.62, RW - 0.24, 0.62,
        [('主档 ＋ 标签 ＋ 兴趣偏好 ＋ 行为摘要（90 天）＋ 会议 /', {'sz': 6.0, 'c': SUB}),
         ('拜访 / 积分 / 工单交互史', {'sz': 6.0, 'c': SUB}),
         ('读模型：Redis 热点 TTL 5min · Cosmos 全量 · 10 万 HCP / 500 QPS', {'sz': 5.9, 'c': MUT})],
        leading=1.22)

    rect(sl, RX, 4.38, RW, 0.92, 'F6F8FA', 'CBD5E1', 1.2, 0.06)
    txt(sl, RX + 0.12, 4.45, RW - 0.24, 0.14, '数据治理', sz=8, b=True, c=SLA)
    txt(sl, RX + 0.12, 4.62, RW - 0.24, 0.62,
        [('分类分级 P1–P4 · 列级血缘（源 → Gold → 看板）', {'sz': 6.0, 'c': SUB}),
         ('质量规则：OneID 唯一性 · 积分账本平衡 · Freshness', {'sz': 6.0, 'c': SUB}),
         ('第三方 HCP 参考库入库前法务评估', {'sz': 5.9, 'c': MUT})], leading=1.22)

    # 中→右：Silver 层 Consent 过滤已在上文 Silver 阶段标注（★Consent 过滤），不画跨列箭头

    # ---- 底部 ----
    rect(sl, 0.22, 5.46, 12.89, 0.44, L_BLUE, B_BLUE, 1.1, 0.06)
    txt(sl, 0.40, 5.46, 12.5, 0.44,
        [dict(runs=[('关键指标：', {'b': True, 'c': '1D4ED8', 'sz': 6.8}),
                    ('OneID 误并率 < 0.5%（阈值保守 ＋ 人工审核）· Consent 撤回生效 ≤ 24h · C360 读 p95 ≤ 100ms · 行为流峰值 5k events/s · Silver 层即执行 Consent 过滤（无 profiling 同意不进画像）',
                     {'c': SUB, 'sz': 6.8})])], anchor='m')
    rect(sl, 0.22, 5.98, 12.89, 0.44, 'FFF7ED', 'FDBA74', 1.1, 0.06)
    txt(sl, 0.40, 5.98, 12.5, 0.44,
        [dict(runs=[('合规映射：', {'b': True, 'c': 'B45309', 'sz': 6.8}),
                    ('PIPL（来源合法 · 最小必要 · 出境评估）｜ 网安法 / 数安法（境内驻留 · 分类分级）｜ 处方药信息（HCP 认证 ＋ 内容受众标签 ＋ 检索行级过滤）',
                     {'c': SUB, 'sz': 6.8})])], anchor='m')
    # 图例
    leg = [('Medallion 分层', SLA), ('OneID 匹配', GRN), ('Consent / 合规', GRN),
           ('标签 / C360', AZ), ('治理', PUR)]
    xx = 0.40
    for t, c in leg:
        oval(sl, xx, 6.60, 0.09, 0.09, c)
        txt(sl, xx + 0.14, 6.55, len(t) * 0.10 + 0.30, 0.16, t, sz=6.4, c=MUT)
        xx += 0.32 + len(t) * 0.10 + 0.50


# ================================================= SLIDE 4 · AI 平台治理 ===
def slide_ai(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _cur['i'] = 4
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor.from_string(BG)
    title_bar(sl, 'AI 平台架构 · 模型网关 · RAG · 医药合规',
              'AI First：可评测 · 可审计 · 可切换 · 可计量\nAI 问答 · 推荐 · 客服 · 合规审核 · 运营分析', PUR)

    # ---- 左：model-gateway ----
    LX = 0.22; LW = 4.10
    rect(sl, LX, 0.62, LW, 0.34, PUR, None, 0, 0.05)
    txt(sl, LX, 0.62, LW, 0.34, 'model-gateway · 统一模型接入（一切调用必经）', sz=7.8, b=True, c='FFFFFF',
        align='c', anchor='m')
    routes = [
        ('非敏场景', '→ Azure OpenAI（全球区 East Asia）', L_BLUE, B_BLUE, AZD),
        ('含个人数据', '→ 境内驻留模型（国产备案 · AKS 推理）', L_GRN, B_GRN, GRND),
        ('简单任务', '→ mini 档小模型（分类 / 审核初筛）', L_AMB, B_AMB, AMBD),
    ]
    yy = 1.04
    for t, s, lf, bl, tc in routes:
        rect(sl, LX, yy, LW, 0.36, lf, bl, 1.1, 0.05)
        txt(sl, LX + 0.10, yy, LW - 0.20, 0.36,
            [dict(runs=[(t + '  ', {'b': True, 'c': tc, 'sz': 6.8}), (s, {'c': SUB, 'sz': 6.4})])],
            anchor='m')
        yy += 0.44
    caps = ['配额 · 成本标签（hub / tenant / 场景）', 'Prompt / 响应全量审计（Cosmos 留痕）',
            '降级链：重试 → 备用模型 → 兜底话术', '语义缓存：高频学术问答命中直返']
    for i, s in enumerate(caps):
        cx = LX + (i % 2) * 2.07; cy = 2.40 + (i // 2) * 0.42
        chip(sl, cx, cy, 1.99, 0.36, s, None, 'FFFFFF', B_PUR, 1.0, 0.04, tc=SUB, tsz=5.9)
    rect(sl, LX, 3.28, LW, 0.66, 'FFF7ED', 'FDBA74', 1.2, 0.06)
    txt(sl, LX + 0.12, 3.28, LW - 0.24, 0.66,
        [('模型路线 C（推荐）', {'b': True, 'sz': 7.2, 'c': 'B45309'}),
         ('AOAI ＋ 驻留模型混合 · 网关按场景路由，规避中国区不可用 / 出境风险；网关抽象保证可切换',
          {'sz': 6.0, 'c': AMBD})], anchor='m', leading=1.2)
    rect(sl, LX, 4.02, LW, 0.56, 'F6F8FA', 'CBD5E1', 1.1, 0.06)
    txt(sl, LX + 0.12, 4.02, LW - 0.24, 0.56,
        [('工具调用：Copilot 仅连只读内部 API（经 APIM）', {'sz': 6.2, 'c': SUB, 'b': True}),
         ('运营 NL2SQL 限白名单视图（受控语义层）', {'sz': 6.0, 'c': MUT})], anchor='m', leading=1.2)

    # ---- 中：RAG 管道 ----
    CX = 4.46; CW = 4.44
    rect(sl, CX, 0.62, CW, 0.34, '5B21B6', None, 0, 0.05)
    txt(sl, CX, 0.62, CW, 0.34, 'RAG 检索增强管道（rag-service）', sz=7.8, b=True, c='FFFFFF',
        align='c', anchor='m')
    rag = [
        ('内容库', 'DAM / CMS 已审核产品与学术内容', SLA, 'F6F8FA', 'CBD5E1'),
        ('解析', 'AI Document Intelligence（PDF / 视频 / 图文）', CYN, L_CYN, B_CYNB),
        ('结构化分块', '保留医学元数据：适应症 · 证据等级 · 有效期', CYN, L_CYN, B_CYNB),
        ('嵌入 ＋ 索引', '双语嵌入 → AI Search（向量＋BM25＋语义排序）', AZ, L_BLUE, B_BLUE),
        ('检索', 'top-k ＋ 引用定位 · 行级过滤（tenant ＋ 角色）', AZ, L_BLUE, B_BLUE),
        ('生成', 'LLM（system prompt 约束 · 引用必附 · 免责声明）', PUR, L_PUR, B_PUR),
    ]
    yy = 1.04
    for i, (t, s, col, lf, bl) in enumerate(rag):
        rect(sl, CX, yy, CW, 0.44, lf, bl, 1.1, 0.05)
        rect(sl, CX + 0.015, yy + 0.05, 0.04, 0.34, col, rad=0.02)
        txt(sl, CX + 0.12, yy + 0.04, 1.05, 0.36, t, sz=7.2, b=True, c=INK, anchor='m')
        txt(sl, CX + 1.18, yy, CW - 1.30, 0.44, s, sz=6.0, c=SUB, anchor='m', leading=1.12)
        if i < 5:
            seg(sl, CX + CW / 2, yy + 0.44, CX + CW / 2, yy + 0.52, MUT, 1.4)
        yy += 0.52
    rect(sl, CX, 4.20, CW, 0.38, L_ROS, ROS, 1.1, 0.05)
    txt(sl, CX + 0.10, 4.20, CW - 0.20, 0.38,
        [dict(runs=[('检索安全：', {'b': True, 'c': ROS, 'sz': 6.4}),
                    ('处方药内容行级过滤 → 仅专业人士渠道可见', {'c': SUB, 'sz': 6.4})])], anchor='m')

    # ---- 右：治理 ----
    RX = 9.04; RW = 4.07
    rect(sl, RX, 0.62, RW, 0.34, '9F1239', None, 0, 0.05)
    txt(sl, RX, 0.62, RW, 0.34, 'AI 治理 · 医药合规', sz=7.8, b=True, c='FFFFFF', align='c', anchor='m')
    rect(sl, RX, 1.04, RW, 0.62, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, RX + 0.12, 1.10, RW - 0.24, 0.52,
        [('两步审核工作流', {'b': True, 'sz': 7.2, 'c': INK}),
         ('AI 预审（review-service）→ 医学 / 合规人审（workflow）', {'sz': 6.0, 'c': SUB}),
         ('→ 发布 · 审核结果回流优化分类器', {'sz': 6.0, 'c': SUB})], leading=1.18)
    rect(sl, RX, 1.74, RW, 0.94, L_ROS, ROS, 1.3, 0.06, dash=DASH)
    txt(sl, RX + 0.12, 1.80, RW - 0.24, 0.84,
        [('AE / PQ 出站护栏 ★（药物警戒）', {'b': True, 'sz': 7.2, 'c': ROS}),
         ('识别不良事件 / 产品质量投诉 → 中断个性化回复 →', {'sz': 6.0, 'c': SUB}),
         ('引导话术 → 自动创建 PV 工单（按 SOP 时效）→', {'sz': 6.0, 'c': SUB}),
         ('会话全量留痕可追溯', {'sz': 6.0, 'c': SUB})], leading=1.18)
    rect(sl, RX, 2.76, RW, 0.62, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, RX + 0.12, 2.82, RW - 0.24, 0.52,
        [('红队与评测', {'b': True, 'sz': 7.2, 'c': INK}),
         ('golden set 回归 ＋ prompt 注入 / 越狱用例库', {'sz': 6.0, 'c': SUB}),
         ('Prompt Flow 评估入 CI · 不通过不上线', {'sz': 6.0, 'c': SUB})], leading=1.18)
    rect(sl, RX, 3.46, RW, 1.12, 'FFFFFF', BORDER, 1.2, 0.06)
    txt(sl, RX + 0.12, 3.52, RW - 0.24, 0.14, 'Copilot 场景（数据最小化）', sz=7.4, b=True, c=INK)
    cps = [('HCP 学术助手', '小程序 · 已审核内容库 RAG'),
           ('销售 Copilot', '企微侧边栏 · 拜访总结 / 客户简报'),
           ('运营 Copilot', 'NL2SQL 白名单视图 · 指标问答'),
           ('服务 Copilot', '工单诊断 · 备件推荐（IoT ＋ KB）')]
    for i, (t, s) in enumerate(cps):
        chip(sl, RX + 0.12, 3.70 + i * 0.215, 1.28, 0.19, t, None, L_PUR, B_PUR, 1.0, 0.04,
             tc=PURD, tsz=5.8)
        txt(sl, RX + 1.50, 3.70 + i * 0.215, RW - 1.62, 0.19, s, sz=5.9, c=MUT, anchor='m')
    rect(sl, RX, 4.66, RW, 0.44, 'F6F8FA', 'CBD5E1', 1.1, 0.06)
    txt(sl, RX + 0.12, 4.66, RW - 0.24, 0.44,
        [('医学边界：system prompt 禁诊疗建议 · 高危转人工', {'sz': 6.0, 'c': SUB, 'b': True}),
         ('输出强制附引用来源与免责声明', {'sz': 6.0, 'c': MUT})], anchor='m', leading=1.15)

    # gateway -> RAG / 治理 关联（列间窄缝，不加文字标注）
    seg(sl, LX + LW, 1.90, CX, 1.90, PUR, 1.4, DASH)
    seg(sl, CX + CW, 2.20, RX, 2.20, ROS, 1.4, DASH)

    # ---- 底部 ----
    rect(sl, 0.22, 5.28, 12.89, 0.44, L_PUR, B_PUR, 1.1, 0.06)
    txt(sl, 0.40, 5.28, 12.5, 0.44,
        [dict(runs=[('AI 能力（archi.md §10）：', {'b': True, 'c': PURD, 'sz': 6.8}),
                    ('AI 问答 · AI 推荐 · AI 客服 · AI 合规审核 · AI 运营分析　｜　统一经 model-gateway：可评测（评估集入 CI）· 可审计（全量留痕）· 可切换（路由抽象）· 可计量（token 打标）',
                     {'c': SUB, 'sz': 6.8})])], anchor='m')
    rect(sl, 0.22, 5.80, 12.89, 0.44, 'FFF7ED', 'FDBA74', 1.1, 0.06)
    txt(sl, 0.40, 5.80, 12.5, 0.44,
        [dict(runs=[('风险与缓解：', {'b': True, 'c': 'B45309', 'sz': 6.8}),
                    ('R1 Azure OpenAI 中国区不可用 / 出境未获批 → 网关抽象 ＋ 国产模型备选（Phase 1 合规预沟通）· R2 微信医疗类目资质 → 提前申报 ＋ H5 降级 · R3 AI 幻觉 → 引用必附 ＋ 两步审核 ＋ 红队回归',
                     {'c': SUB, 'sz': 6.8})])], anchor='m')
    txt(sl, 0.22, 6.42, 12.9, 0.5,
        [('落地节奏：Phase 1 统一底座（OneID / CMS / Event / HCP Pool）→ Phase 2 能力复用（微前端 / Capability Hub / 插件 / C360）→ Phase 3 AI 赋能（AI Search / 推荐 / Copilot / AI 合规）',
          {'sz': 6.8, 'c': MUT})], align='c', anchor='m')


# ------------------------------------------------------------------- build --
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = 'One-Abbott 共享应用平台 · 技术架构图集'
    prs.core_properties.author = 'One-Abbott Platform'
    slide_cover(prs)
    slide_overall(prs)
    slide_mfe(prs)
    slide_hcp(prs)
    slide_ai(prs)
    import os as _os
    out = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), OUT_NAME)
    prs.save(out)
    with open('/tmp/arch_ops.json', 'w', encoding='utf-8') as f:
        import json
        json.dump(OPS, f, ensure_ascii=False)
    print('saved', OUT_NAME, '| shapes ops:', len(OPS))


if __name__ == '__main__':
    build()
